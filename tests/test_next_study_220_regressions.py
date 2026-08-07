"""Regression probes for defects that invalidated the v0.13.0 instrument."""

import copy
from io import BytesIO
from pathlib import Path
import re

from openpyxl import load_workbook
from pptx import Presentation

from bench.next_study_grader_audit import build_positive_evidence, _rebuild
from bench.next_study_review import digest_review_artifact, review_packet
from domains.office_demo.generators_v2 import generate_instance
from domains.office_demo.outcome_oracle_v2 import derive_outcome
from domains.office_demo.reviewed_grader_v2 import build_grader


def _case(family, index=0):
    instance = generate_instance("adversarial", family, index)
    packet = review_packet(instance)
    effects = derive_outcome(
        packet["family"], packet["prompt"], packet["subepisode_prompts"],
        packet["initial_state"], packet["today"],
    )
    outcome = {
        "instance_id": instance["content"]["id"],
        "content_sha256": instance["content_sha256"],
        "review_packet_sha256": digest_review_artifact(packet),
        "prompt_valid": True,
        "outcome": effects,
        "accepted_alternatives": [],
        "review_resolution": "deterministic_regression_fixture",
    }
    return instance, packet, outcome


def _baseline(tmp_path, family, index=0):
    instance, packet, outcome = _case(family, index)
    evidence = build_positive_evidence(packet, outcome, tmp_path / family)
    assert build_grader(packet, outcome).grade_evidence(evidence).strict_success
    return instance, packet, outcome, evidence


def test_offsite_issued_rank_agrees_with_visible_date_order():
    for index in range(4):
        instance = generate_instance("adversarial", "multi_offsite", index)
        emails = {
            item["id"]: item for item in instance["content"]["initial_state"]["emails"]
        }
        candidates = [item for item in emails if item.startswith("offsite-final-")]
        by_date = sorted(candidates, key=lambda item: emails[item]["date"])
        index_body = emails["offsite-index"]["body"]
        by_rank = sorted(
            candidates,
            key=lambda item: int(index_body.split("id=" + item + ",issued_rank=")[1].split(",", 1)[0]),
        )
        assert by_date == by_rank


def test_preference_title_prefix_is_in_saved_public_facts():
    for index in range(4):
        instance = generate_instance("adversarial", "preference_learning", index)
        store = instance["content"]["ordered_subepisodes"][0]["prompt"]
        facts = store.split("applicable facts:", 1)[1]
        assert "title_prefix=" in facts


def test_reminder_prompt_does_not_print_the_derived_order():
    instance = generate_instance("adversarial", "remind_msg", 0)
    prompt = instance["content"]["prompt"]
    assert "Required checklist mentions" not in prompt
    assert "resulting full ordered ID list" in prompt


def test_receipt_prompt_pins_dollar_conversion_and_grader_enforces_it(tmp_path):
    for index in range(4):
        instance = generate_instance("adversarial", "xlsx_from_email", index)
        assert "USD dollar values" in instance["content"]["prompt"]
        assert "amount_cents=N" in instance["content"]["prompt"]
        assert "N/100 dollars" in instance["content"]["prompt"]

    _instance, packet, outcome, evidence = _baseline(tmp_path, "xlsx_from_email")
    effect = next(item for item in outcome["outcome"] if item["type"] == "spreadsheet_created")
    artifacts = dict(evidence.artifact_map())
    workbook = load_workbook(BytesIO(artifacts[effect["filename"]]))
    sheet = workbook.active
    assert float(sheet.cell(row=2, column=3).value) == (
        effect["ordered_rows_cents"][0][2] / 100
    )
    for row, wanted in enumerate(effect["ordered_rows_cents"], start=2):
        sheet.cell(row=row, column=3).value = wanted[2]
    output = BytesIO()
    workbook.save(output)
    workbook.close()
    artifacts[effect["filename"]] = output.getvalue()
    assert not build_grader(packet, outcome).grade_evidence(
        _rebuild(evidence, artifacts=sorted(artifacts.items()))
    ).strict_success


def test_free_slot_prompt_defines_policy_precedence():
    for index in range(4):
        prompt = generate_instance("adversarial", "cal_freeslot", index)["content"]["prompt"]
        assert "Apply only the named policy" in prompt
        assert "closest_to_preferred selects" in prompt
        assert "equal-distance tie toward the earlier slot" in prompt


def test_reply_prompt_defines_policy_and_exact_field_copy():
    for index in range(4):
        prompt = generate_instance("adversarial", "email_reply", index)["content"]["prompt"]
        assert "decision_key exactly equals" in prompt
        assert "copying all three field values exactly as shown" in prompt
        assert "Do not reply to any other sender or create any reminder" in prompt


def test_reply_rejects_reformatted_required_value_and_extra_business_action(tmp_path):
    _instance, packet, outcome, evidence = _baseline(tmp_path, "email_reply")
    grader = build_grader(packet, outcome)
    effect = next(item for item in outcome["outcome"] if item["type"] == "email_sent")
    state = copy.deepcopy(evidence.state)
    iso_date = effect["required_mentions"][1]
    state["sent_emails"][-1]["body"] = state["sent_emails"][-1]["body"].replace(
        iso_date, "February 5, 2028"
    )
    assert not grader.grade_evidence(_rebuild(evidence, state=state)).strict_success

    state = copy.deepcopy(evidence.state)
    state["reminders"].append({"text": "unrequested", "date": iso_date, "time": "09:00"})
    actions = list(evidence.actions) + [{
        "tool": "set_reminder",
        "args": {"text": "unrequested", "date": iso_date, "time": "09:00"},
        "ok": True,
        "result": "created",
    }]
    assert not grader.grade_evidence(
        _rebuild(evidence, state=state, actions=actions)
    ).strict_success


def test_reminder_dependencies_are_due_date_coherent_and_outputs_are_exact(tmp_path):
    for index in range(4):
        instance = generate_instance("adversarial", "remind_msg", index)
        prompt = instance["content"]["prompt"]
        records = {}
        record_text = prompt.split("Action items: ", 1)[1].split(". Order them", 1)[0]
        for value in record_text.split(" | "):
            fields = dict(part.split("=", 1) for part in value.split(","))
            records[fields["id"]] = fields
        for record in records.values():
            if record["depends_on"] != "none":
                assert records[record["depends_on"]]["due"] <= record["due"]
        assert "breaking multiple-eligible ties by larger priority" in prompt
        assert "first ordered item's due date" in prompt

    _instance, packet, outcome, evidence = _baseline(tmp_path, "remind_msg")
    grader = build_grader(packet, outcome)
    state = copy.deepcopy(evidence.state)
    state["reminders"][-1]["text"] += "; checkpoint-99"
    assert not grader.grade_evidence(_rebuild(evidence, state=state)).strict_success
    state = copy.deepcopy(evidence.state)
    state["messages"][-1]["text"] += "; checkpoint-99"
    assert not grader.grade_evidence(_rebuild(evidence, state=state)).strict_success
    state = copy.deepcopy(evidence.state)
    state["messages"][-1]["text"] = state["messages"][-1]["text"].replace(
        next(item["deadline"] for item in outcome["outcome"] if item["type"] == "message_sent"),
        "the deadline",
    )
    assert not grader.grade_evidence(_rebuild(evidence, state=state)).strict_success


def test_offsite_prompt_defines_rank_direction():
    for index in range(4):
        prompt = generate_instance("adversarial", "multi_offsite", index)["content"]["prompt"]
        assert "largest issued_rank value" in prompt
        assert "largest approval_rank value" in prompt
        assert "largest consensus value" in prompt


def test_missing_inbox_listing_is_rejected(tmp_path):
    _instance, packet, outcome, evidence = _baseline(tmp_path, "email_reply")
    changed = _rebuild(
        evidence,
        actions=[item for item in evidence.actions if item["tool"] != "list_emails"],
    )
    assert not build_grader(packet, outcome).grade_evidence(changed).strict_success


def test_negated_affirmative_intents_are_rejected(tmp_path):
    for family, text in (
        ("email_reply", "I cannot confirm and will not attend."),
        ("remind_msg", "This will not be complete by the deadline."),
    ):
        _instance, packet, outcome, evidence = _baseline(tmp_path, family)
        state = copy.deepcopy(evidence.state)
        if family == "email_reply":
            effect = next(item for item in outcome["outcome"] if item["type"] == "email_sent")
            state["sent_emails"][-1]["body"] = text + " " + " ".join(
                effect.get("required_mentions", [])
            )
        else:
            effect = next(item for item in outcome["outcome"] if item["type"] == "message_sent")
            state["messages"][-1]["text"] = "; ".join(effect["required_mentions"]) + "; " + text
        assert not build_grader(packet, outcome).grade_evidence(
            _rebuild(evidence, state=state)
        ).strict_success


def test_calendar_brief_rejects_forbidden_date(tmp_path):
    _instance, packet, outcome, evidence = _baseline(tmp_path, "cal_brief")
    state = copy.deepcopy(evidence.state)
    state["messages"][0]["text"] += "; 2099-12-31"
    assert not build_grader(packet, outcome).grade_evidence(
        _rebuild(evidence, state=state)
    ).strict_success


def test_calendar_brief_rejects_excluded_title_and_pins_entry_format(tmp_path):
    instance, packet, outcome, evidence = _baseline(tmp_path, "cal_brief", index=1)
    assert "formatting every entry exactly as '<title> at <HH:MM>'" in packet["prompt"]
    excluded = next(
        event["title"] for event in instance["content"]["initial_state"]["events"]
        if not event["title"].startswith("Priority:")
    )
    state = copy.deepcopy(evidence.state)
    state["messages"][0]["text"] += "; " + excluded
    assert not build_grader(packet, outcome).grade_evidence(
        _rebuild(evidence, state=state)
    ).strict_success


def test_basic_families_do_not_require_nonexistent_source_reads(tmp_path):
    for family in ("pptx_basic", "xlsx_basic"):
        _instance, packet, outcome, evidence = _baseline(tmp_path, family)
        assert not any(action["tool"] == "read_email" for action in evidence.actions)
        assert build_grader(packet, outcome).grade_evidence(evidence).strict_success


def test_presentation_numeric_superstring_is_rejected(tmp_path):
    _instance, packet, outcome, evidence = _baseline(tmp_path, "pptx_from_email")
    effect = next(item for item in outcome["outcome"] if item["type"] == "presentation_created")
    artifacts = dict(evidence.artifact_map())
    deck = Presentation(Path(tmp_path / "pptx_from_email" / "artifacts" / effect["filename"]))
    changed = False
    for slide in deck.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text.strip().isdigit() and not changed:
                    paragraph.text = paragraph.text.strip() + "0"
                    changed = True
    output = tmp_path / "numeric-superstring.pptx"
    deck.save(output)
    artifacts[effect["filename"]] = output.read_bytes()
    assert changed
    assert not build_grader(packet, outcome).grade_evidence(
        _rebuild(evidence, artifacts=sorted(artifacts.items()))
    ).strict_success


def test_scratch_memory_is_neutral_but_requested_memory_is_exact(tmp_path):
    _instance, packet, outcome, evidence = _baseline(tmp_path, "cal_add")
    scratch = "temporary internal scratch note"
    changed = _rebuild(
        evidence,
        memory=list(evidence.memory) + [scratch],
        actions=list(evidence.actions) + [{
            "tool": "save_memory", "args": {"fact": scratch}, "ok": True,
            "result": "saved",
        }],
    )
    assert build_grader(packet, outcome).grade_evidence(changed).strict_success

    _instance, packet, outcome, evidence = _baseline(tmp_path, "preference_learning")
    changed = _rebuild(
        evidence,
        memory=list(evidence.memory) + ["unrequested extra memory"],
        actions=list(evidence.actions) + [{
            "tool": "save_memory", "args": {"fact": "unrequested extra memory"},
            "ok": True, "result": "saved",
        }],
    )
    assert not build_grader(packet, outcome).grade_evidence(changed).strict_success


def test_reported_prompt_grader_blockers_reproduce(tmp_path):
    """Pin the ten defect classes that terminally block office-generators/2.1.2."""

    # The prompt permits a formula, but the evaluator accepts only SUM(range).
    _instance, packet, outcome, evidence = _baseline(tmp_path, "xlsx_basic")
    effect = next(item for item in outcome["outcome"] if item["type"] == "spreadsheet_created")
    artifacts = dict(evidence.artifact_map())
    workbook = load_workbook(BytesIO(artifacts[effect["filename"]]))
    sheet = workbook.active
    last_data_row = sheet.max_row - 1
    total_column = sheet.max_column
    sheet.cell(row=sheet.max_row, column=total_column).value = "=" + "+".join(
        "%s%d" % (sheet.cell(row=1, column=total_column).column_letter, row)
        for row in range(2, last_data_row + 1)
    )
    output = BytesIO()
    workbook.save(output)
    workbook.close()
    artifacts[effect["filename"]] = output.getvalue()
    assert not build_grader(packet, outcome).grade_evidence(
        _rebuild(evidence, artifacts=sorted(artifacts.items()))
    ).strict_success

    # A natural confirmation is not in the hidden reviewed-grader allowlist.
    _instance, packet, outcome, evidence = _baseline(tmp_path, "email_reply")
    effect = next(item for item in outcome["outcome"] if item["type"] == "email_sent")
    state = copy.deepcopy(evidence.state)
    state["sent_emails"][-1]["body"] = "I confirm my attendance. " + "; ".join(
        effect["required_mentions"]
    )
    assert not build_grader(packet, outcome).grade_evidence(
        _rebuild(evidence, state=state)
    ).strict_success

    # highest_priority_feasible never defines whether larger numbers rank higher.
    priority_cases = [
        generate_instance("development", "cal_add", index)
        for index in range(8)
    ]
    priority = next(
        item for item in priority_cases
        if item["content"]["structure"]["decision_policy"]
        == "highest_priority_feasible"
    )
    assert "largest priority value" not in priority["content"]["prompt"]

    # "titled for" permits decoration, but the live grader requires exact equality.
    _instance, packet, outcome, evidence = _baseline(tmp_path, "multi_offsite")
    state = copy.deepcopy(evidence.state)
    state["events"][-1]["title"] = "Offsite: " + state["events"][-1]["title"]
    assert not build_grader(packet, outcome).grade_evidence(
        _rebuild(evidence, state=state)
    ).strict_success

    # A natural deadline commitment outside the hidden allowlist also fails.
    _instance, packet, outcome, evidence = _baseline(tmp_path, "remind_msg")
    effect = next(item for item in outcome["outcome"] if item["type"] == "message_sent")
    state = copy.deepcopy(evidence.state)
    state["messages"][-1]["text"] = "%s; I'll have everything done by %s." % (
        "; ".join(effect["required_mentions"]), effect["deadline"]
    )
    assert not build_grader(packet, outcome).grade_evidence(
        _rebuild(evidence, state=state)
    ).strict_success

    # The prompt says include the fact; a faithful label nevertheless fails.
    _instance, packet, outcome, evidence = _baseline(tmp_path, "pptx_basic")
    effect = next(item for item in outcome["outcome"] if item["type"] == "presentation_created")
    artifacts = dict(evidence.artifact_map())
    deck = Presentation(BytesIO(artifacts[effect["filename"]]))
    wanted = str(effect["required_values_by_slide"][1][0])
    changed = False
    for slide in deck.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text.strip() == wanted and not changed:
                    paragraph.text = "Fact: " + wanted
                    changed = True
    output_path = tmp_path / "labelled-fact.pptx"
    deck.save(output_path)
    assert changed
    artifacts[effect["filename"]] = output_path.read_bytes()
    assert not build_grader(packet, outcome).grade_evidence(
        _rebuild(evidence, artifacts=sorted(artifacts.items()))
    ).strict_success

    # The prompt prints pipe-separated facts; the grader silently requires ';'.
    _instance, packet, outcome, evidence = _baseline(tmp_path, "preference_learning")
    changed_memory = list(evidence.memory)
    changed_memory[-1] = changed_memory[-1].replace(";", " | ")
    assert changed_memory[-1] != evidence.memory[-1]
    assert not build_grader(packet, outcome).grade_evidence(
        _rebuild(evidence, memory=changed_memory)
    ).strict_success

    # The same family also has an undisclosed case-sensitive title grammar.
    state = copy.deepcopy(evidence.state)
    state["events"][-1]["title"] = state["events"][-1]["title"].replace(
        " sync with ", " Sync with ", 1
    )
    assert state["events"][-1]["title"] != evidence.state["events"][-1]["title"]
    assert not build_grader(packet, outcome).grade_evidence(
        _rebuild(evidence, state=state)
    ).strict_success

    # One pptx policy is named but has no public definition.
    brief_cases = [
        generate_instance("development", "pptx_basic", index)
        for index in range(8)
    ]
    brief = next(
        item for item in brief_cases
        if item["content"]["structure"]["decision_policy"] == "brief_sequence"
    )
    prompt = brief["content"]["prompt"]
    assert "brief_sequence" in prompt
    assert "brief_sequence selects" not in prompt

    # Enumerated values are silently order-sensitive in the reviewed grader.
    _instance, packet, outcome, evidence = _baseline(tmp_path, "email_reply")
    effect = next(item for item in outcome["outcome"] if item["type"] == "email_sent")
    state = copy.deepcopy(evidence.state)
    state["sent_emails"][-1]["body"] = "I will attend. " + "; ".join(
        reversed(effect["required_mentions"])
    )
    assert not build_grader(packet, outcome).grade_evidence(
        _rebuild(evidence, state=state)
    ).strict_success


def test_successor_live_path_uses_only_reviewed_grader():
    source = (Path(__file__).resolve().parents[1] / "bench" / "next_study_live.py").read_text(
        encoding="utf-8"
    )
    assert "from domains.office_demo.reviewed_grader_v2 import" in source
    assert "domains.office_demo.generated_grader" not in source


def test_cal_add_feasibility_is_currently_vacuous():
    for split, count in (
        ("development", 8), ("calibration", 8), ("validation", 4),
        ("sentinel", 4), ("adversarial", 4),
    ):
        for index in range(count):
            content = generate_instance(split, "cal_add", index)["content"]
            candidates = re.findall(
                r"start=(\d\d:\d\d),duration=(\d+)", content["prompt"]
            )
            assert len(candidates) == 3
            occupied = [
                (event["start"], event["end"])
                for event in content["initial_state"]["events"]
            ]
            for start, duration in candidates:
                hour, minute = map(int, start.split(":"))
                end_minutes = hour * 60 + minute + int(duration)
                end = "%02d:%02d" % divmod(end_minutes, 60)
                assert all(end <= old_start or start >= old_end for old_start, old_end in occupied)
