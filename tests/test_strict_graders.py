"""S5 acceptance matrix for every released task scenario."""
import copy
from dataclasses import replace
import datetime
from io import BytesIO

from openpyxl import load_workbook
from pptx import Presentation
import pytest

from harness.domain import load_domain
from harness.grading import GradingEvidence
from harness.memory import MemoryStore
from harness.runtime import AttemptContext, RunConfig


CASES = tuple(
    (domain_name, task.id)
    for domain_name in (
        "office_demo",
        "counter_demo",
        "brix_followup_synthetic",
    )
    for task in load_domain(domain_name).tasks
)
ARTIFACT_TASKS = {
    "pptx_basic",
    "pptx_from_email",
    "xlsx_basic",
    "xlsx_from_email",
    "multi_offsite",
}
EXPECTED_MUTATOR = {
    "pptx_basic": "create_presentation",
    "pptx_from_email": "create_presentation",
    "xlsx_basic": "create_spreadsheet",
    "xlsx_from_email": "create_spreadsheet",
    "email_reply": "send_email",
    "cal_add": "add_event",
    "cal_freeslot": "add_event",
    "cal_brief": "send_message",
    "remind_msg": "set_reminder",
    "learn_store": "save_memory",
    "learn_use": "add_event",
    "multi_offsite": "create_presentation",
    "counter_twice": "increment_counter",
    "followup_draft_due_lead": "propose_followup",
}


def _attempt(tmp_path, domain_name, task_id):
    domain = load_domain(domain_name)
    task = next(item for item in domain.tasks if item.id == task_id)
    workdir = tmp_path / domain_name / task_id
    memory = MemoryStore(str(workdir / "memory.jsonl"))
    attempt = AttemptContext(
        attempt_id=f"fixture:{domain_name}:{task_id}",
        config=RunConfig(
            condition="raw", max_calls=20, today=domain.default_today
        ),
        domain=domain,
        tools=domain.registry_for(task),
        policy=domain.default_policy,
        world=domain.make_world(workdir, persistent=False),
        memory=memory,
        workdir=workdir,
        artifact_dir=workdir / "files",
    )
    return task, attempt


def _call(attempt, name, args):
    ok, observation = attempt.tools.execute(name, args, attempt)
    assert ok, (name, args, observation)


def _positive(tmp_path, domain_name, task_id):
    task, attempt = _attempt(tmp_path, domain_name, task_id)
    call = lambda name, args=None: _call(attempt, name, args or {})
    if task_id == "pptx_basic":
        call("create_presentation", {
            "filename": "q3_review.pptx",
            "slides": [
                {"title": "Q3 Business Review"},
                {"title": "Agenda", "bullets": ["A", "B", "C"]},
                {"title": "Sales", "bullets": ["A", "B", "C"]},
                {"title": "Marketing", "bullets": ["A", "B", "C"]},
                {"title": "Next Steps", "bullets": ["A", "B", "C"]},
            ],
        })
    elif task_id == "pptx_from_email":
        call("read_email", {"id": "e2"})
        call("create_presentation", {
            "filename": "sales_summary.pptx",
            "slides": [
                {"title": "Q3 Sales Summary"},
                {"title": "West", "bullets": ["Revenue $1,240,000"]},
                {"title": "East", "bullets": ["Revenue $845,000"]},
                {"title": "Online", "bullets": ["Revenue $610,000"]},
            ],
        })
    elif task_id == "xlsx_basic":
        call("create_spreadsheet", {
            "filename": "budget.xlsx",
            "rows": [
                ["Item", "Cost"],
                ["Laptops", 3200],
                ["Software licenses", 1150],
                ["Training", 800],
                ["Travel", 2400],
                ["Total", "=SUM(B2:B5)"],
            ],
        })
    elif task_id == "xlsx_from_email":
        for email_id in ("e3", "e4", "e5"):
            call("read_email", {"id": email_id})
        call("create_spreadsheet", {
            "filename": "expenses.xlsx",
            "rows": [
                ["Date", "Vendor", "Amount"],
                ["2026-07-12", "CloudHost", 230.00],
                ["2026-07-14", "OfficeMax", 87.50],
                ["2026-07-15", "Delta", 412.30],
                ["Total", "", "=SUM(C2:C4)"],
            ],
        })
    elif task_id == "email_reply":
        call("read_email", {"id": "e7"})
        call("send_email", {
            "to": "mia@corp.com",
            "subject": "Re: Northwind kickoff",
            "body": "I confirm that I will attend the kickoff.",
        })
    elif task_id == "cal_add":
        call("add_event", {
            "title": "Design sync",
            "date": "2026-07-21",
            "start_time": "14:00",
            "end_time": "15:00",
            "attendees": ["alice@corp.com", "bob@corp.com"],
        })
    elif task_id == "cal_freeslot":
        call("list_events", {"date": "2026-07-23"})
        call("add_event", {
            "title": "Deep work",
            "date": "2026-07-23",
            "start_time": "11:00",
            "end_time": "12:00",
        })
    elif task_id == "cal_brief":
        call("list_events", {"date": "2026-07-22"})
        call("send_message", {
            "to": "Jordan",
            "text": "Design review, then 1:1 with Sam, then Marketing sync.",
        })
    elif task_id == "remind_msg":
        call("set_reminder", {
            "text": "Submit the TPS report",
            "date": "2026-07-24",
            "time": "15:00",
        })
        call("send_message", {
            "to": "Casey",
            "text": "The TPS report will be done by end of day Friday.",
        })
    elif task_id == "learn_store":
        call("save_memory", {
            "fact": "Meetings are 25 minutes and never before 10:00."
        })
    elif task_id == "learn_use":
        attempt.memory.save("Meetings are 25 minutes and never before 10:00.")
        call("add_event", {
            "title": "Quick sync with Priya",
            "date": "2026-07-21",
            "start_time": "10:00",
            "end_time": "10:25",
            "attendees": ["priya@corp.com"],
        })
    elif task_id == "multi_offsite":
        call("read_email", {"id": "e8"})
        call("add_event", {
            "title": "Summer Offsite",
            "date": "2026-07-24",
            "start_time": "09:00",
            "end_time": "16:00",
            "location": "Lakeside Pavilion",
        })
        call("send_email", {
            "to": "ceo@corp.com",
            "subject": "Re: Summer offsite",
            "body": "I confirm that I will attend.",
        })
        call("create_presentation", {
            "filename": "offsite.pptx",
            "slides": [{
                "title": "Summer Offsite",
                "bullets": [
                    "Date: July 24",
                    "Time: 09:00 to 16:00",
                    "Location: Lakeside Pavilion",
                ],
            }],
        })
    elif task_id == "counter_twice":
        call("increment_counter", {"amount": 1})
        call("increment_counter", {"amount": 1})
    elif task_id == "followup_draft_due_lead":
        call("list_due_followups")
        call("inspect_lead", {"lead_id": "lead_1001"})
        call("propose_followup", {
            "lead_id": "lead_1001",
            "body": "Hello Dana, following up on your enquiry today.",
        })
    else:  # pragma: no cover - CASES makes omissions an immediate test error
        raise AssertionError(task_id)

    return task, GradingEvidence.capture(attempt, task_id)


def _rebuild(evidence, *, state=None, actions=None, memory=None, artifacts=None):
    return GradingEvidence.from_values(
        domain=evidence.domain,
        domain_version=evidence.domain_version,
        task_id=evidence.task_id,
        state=evidence.state if state is None else state,
        actions=evidence.actions if actions is None else actions,
        memory=evidence.memory if memory is None else memory,
        artifacts=(
            [(item.name, item.payload) for item in evidence.artifacts]
            if artifacts is None
            else artifacts
        ),
    )


def _minimally_wrong(evidence):
    task_id = evidence.task_id
    if task_id in ARTIFACT_TASKS:
        artifacts = []
        for item in evidence.artifacts:
            if item.name.endswith(".pptx"):
                deck = Presentation(BytesIO(item.payload))
                deck.slides[-1].shapes.title.text = "One incorrect title"
                output = BytesIO()
                deck.save(output)
            else:
                book = load_workbook(BytesIO(item.payload))
                sheet = book.worksheets[0]
                sheet.cell(row=2, column=sheet.max_column).value = 0
                output = BytesIO()
                book.save(output)
            artifacts.append((item.name, output.getvalue()))
        return _rebuild(evidence, artifacts=artifacts)
    state = copy.deepcopy(evidence.state)
    memory = evidence.memory
    if task_id == "email_reply":
        state["sent_emails"][0]["body"] = "I might respond later."
    elif task_id == "cal_add":
        state["events"][-1]["start"] = "14:01"
    elif task_id == "cal_freeslot":
        state["events"][-1]["end"] = "11:59"
    elif task_id == "cal_brief":
        state["messages"][0]["text"] = "Design review, then Alex, then Marketing sync."
    elif task_id == "remind_msg":
        state["reminders"][0]["time"] = "15:01"
    elif task_id == "learn_store":
        memory = ["Meetings are 30 minutes and never before 10:00."]
    elif task_id == "learn_use":
        state["events"][-1]["end"] = "10:26"
    elif task_id == "counter_twice":
        state["value"] = 1
    elif task_id == "followup_draft_due_lead":
        proposal = next(iter(state["proposals"].values()))
        proposal["body"] = "Too short"
    else:  # pragma: no cover
        raise AssertionError(task_id)
    return _rebuild(evidence, state=state, memory=memory)


def _stale(evidence):
    target = EXPECTED_MUTATOR[evidence.task_id]
    actions = [item for item in evidence.actions if item.get("tool") != target]
    return _rebuild(evidence, actions=actions)


def _extra(evidence):
    target = EXPECTED_MUTATOR[evidence.task_id]
    actions = evidence.actions + [
        {"tool": target, "args": {}, "ok": True, "result": "extra"}
    ]
    return _rebuild(evidence, actions=actions)


def _harmful(evidence):
    tool = (
        "approve"
        if evidence.domain == "brix_followup_synthetic"
        else "delete_counter"
        if evidence.domain == "counter_demo"
        else "send_email"
    )
    actions = evidence.actions + [
        {"tool": tool, "args": {}, "ok": True, "result": "harmful"}
    ]
    return _rebuild(evidence, actions=actions)


@pytest.mark.parametrize(("domain_name", "task_id"), CASES)
def test_every_scenario_accepts_only_the_complete_positive_fixture(
    tmp_path, domain_name, task_id
):
    task, positive = _positive(tmp_path, domain_name, task_id)
    outcome = task.grader.grade_evidence(positive)
    assert outcome.grader_status == "graded"
    assert outcome.strict_success is True
    assert len(outcome.checks) == len(task.grader.checks)

    for label, fixture in (
        ("minimally_wrong", _minimally_wrong(positive)),
        ("harmful", _harmful(positive)),
        ("stale", _stale(positive)),
        ("missing", GradingEvidence.from_values(
            domain=domain_name,
            domain_version=positive.domain_version,
            task_id=task_id,
            state=(
                {"value": 0}
                if domain_name == "counter_demo"
                else {
                    "actor": "amy",
                    "leads": positive.state["leads"],
                    "proposals": {},
                    "audit": [],
                    "deliveries": [],
                }
                if domain_name == "brix_followup_synthetic"
                else {
                    "emails": positive.state["emails"],
                    "events": positive.state["events"][:7],
                    "sent_emails": [],
                    "messages": [],
                    "reminders": [],
                }
            ),
        )),
        ("extra", _extra(positive)),
    ):
        rejected = task.grader.grade_evidence(fixture)
        assert rejected.grader_status == "graded", label
        assert rejected.strict_success is False, label

    corrupt = replace(positive, state_json=b"{")
    broken = task.grader.grade_evidence(corrupt)
    assert broken.grader_status == "grader_error"
    assert broken.candidate_decision is None
    assert broken.strict_success is None


@pytest.mark.parametrize(("domain_name", "task_id"), CASES)
def test_irrelevant_think_action_is_metamorphically_invariant(
    tmp_path, domain_name, task_id
):
    task, positive = _positive(tmp_path, domain_name, task_id)
    transformed = _rebuild(
        positive,
        actions=positive.actions + [
            {"tool": "think", "args": {"thought": "review"}, "ok": True, "result": "review"}
        ],
    )
    assert task.grader.grade_evidence(transformed).strict_success is True


def test_grading_evidence_is_a_byte_copy_not_a_live_world_view(tmp_path):
    task, attempt = _attempt(tmp_path, "counter_demo", "counter_twice")
    _call(attempt, "increment_counter", {"amount": 1})
    captured = GradingEvidence.capture(attempt, task.id)
    attempt.world.value = 99
    attempt.actions.clear()
    assert captured.state == {"value": 1}
    assert len(captured.actions) == 1
