"""Agent Lab — a local web console for the per-model agents.

Pick a model, type a task, press Run, and watch the loop work: the plan, each
model call streaming token by token, every tool call with the arguments the
harness actually sent, and the selected domain's state/files/memory updating
as it changes.

    python -m webui.server            then open http://127.0.0.1:8765

Binds loopback only. One run at a time, in a subprocess (webui/runner.py), so
Stop can terminate that isolated active run.
"""
import http.server
import json
import mimetypes
import os
import queue
from collections import deque
import subprocess
import sys
import threading
import time
import urllib.parse
import webbrowser
import zipfile

import requests

HERE = os.path.dirname(os.path.abspath(__file__))
PROJECT = os.path.dirname(HERE)
if PROJECT not in sys.path:
    sys.path.insert(0, PROJECT)
STATIC = os.path.join(HERE, "static")
AGENTS_DIR = os.path.join(PROJECT, "agents")
OLLAMA_URL = "http://127.0.0.1:11434"
DEFAULT_PORT = 8765

from agents._shared.run_agent import validate_config  # noqa: E402
from harness.domain import load_domain  # noqa: E402
from harness import chat  # noqa: E402
from harness import mcp_config  # noqa: E402
from harness.storage import agent_runtime_paths  # noqa: E402
from webui.control import (  # noqa: E402
    ConfirmationLedger,
    EventJournal,
    MAX_ARCHIVE_EXPANDED_BYTES,
    MAX_ARCHIVE_MEMBERS,
    MAX_ARTIFACT_BYTES,
    MAX_LOG_FILE_BYTES,
    MAX_PREVIEW_BYTES,
    MAX_STDERR_LINE,
    MAX_STDERR_LINES,
    ProcessTree,
    RequestError,
    exact_object,
    new_capability,
    portable_leaf,
    read_json_object,
    redact,
    regular_entries_under,
    regular_path_under,
    reset_directory,
    require_bool,
    require_int,
    require_optional_string,
    require_string,
    validate_capability,
    validate_host,
    validate_mutation_origin,
    trusted_directory_under,
    validate_regular_tree_under,
)


REMOVED_RUN_FIELDS = frozenset(
    {
        "root",
        "shell",
        "yolo",
        "with_domain",
        "with_office",
    }
)

# Rough per-size guidance for the picker; the machine, not the harness, decides.
SPEED_HINT = {
    "1b": (
        "1B tier",
        "Smallest configured parameter tier; measure quality and latency "
        "on the target hardware.",
    ),
    "3b": (
        "3B tier",
        "Configured 3B parameter tier; performance is hardware-dependent.",
    ),
    "8b": (
        "8B tier",
        "Configured 8B parameter tier; performance is hardware-dependent.",
    ),
    "14b": (
        "14B tier",
        "Configured 14B parameter tier; verify memory and latency locally.",
    ),
    "32b": (
        "32B tier",
        "Largest configured parameter tier; verify fit and latency locally.",
    ),
}

# ----------------------------------------------------------------- agents ----

def agent_folders():
    if not os.path.isdir(AGENTS_DIR):
        return []
    out = []
    for name in sorted(os.listdir(AGENTS_DIR), key=lambda n: (len(n), n)):
        if name.startswith("_") or name.startswith("."):
            continue
        cfg_path = os.path.join(AGENTS_DIR, name, "config.json")
        if os.path.isfile(cfg_path):
            out.append(name)
    return out


def read_config(agent):
    with open(os.path.join(AGENTS_DIR, agent, "config.json"), encoding="utf-8-sig") as f:
        config = json.load(f)
    validate_config(config)
    return config


def reject_removed_run_fields(body):
    """Fail before process creation when a retired capability is requested."""
    if not isinstance(body, dict):
        raise ValueError("request body must be a JSON object")
    removed = sorted(REMOVED_RUN_FIELDS & set(body))
    if removed:
        raise ValueError(
            "unsupported Agent Lab fields: " + ", ".join(removed)
        )


# Model facts from openrouter.ai, generated into model_catalog.json rather than
# fetched. The product's whole claim is that nothing leaves the machine, so the
# UI must not reach the network to describe a model.
def _load_catalog():
    path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                        "model_catalog.json")
    try:
        with open(path, encoding="utf-8") as handle:
            return json.load(handle).get("models", {})
    except (OSError, ValueError):
        return {}


CATALOG = _load_catalog()


def catalog_for(tag):
    """Exact tag first, then the bare family, so llama3.1 matches llama3.1:8b."""
    if tag in CATALOG:
        return CATALOG[tag]
    base = tag.split(":")[0]
    for key, value in CATALOG.items():
        if key.split(":")[0] == base:
            return value
    return {}


MCP_MODES = ("draft", "live", "read_only")


def require_mcp_names(value):
    """Connector names from the browser, checked against the registry.

    This is a trust boundary: the names end up as runner argv. Only keys that
    exist in mcp/servers.json get through, so a typo fails here with a clear
    message rather than deep inside a subprocess."""
    if not value:
        return []
    if not isinstance(value, list) or len(value) > 8:
        raise RequestError(400, "mcp must be a list of at most 8 server names")
    known = {name for name, _ in mcp_config.available()}
    unknown = sorted({v for v in value if v not in known})
    if unknown:
        raise RequestError(400, "unknown MCP servers: " + ", ".join(unknown))
    return sorted(set(value))


def require_mcp_mode(value):
    if value is None:
        return None
    if value not in MCP_MODES:
        raise RequestError(400, "mcp_mode must be one of " + ", ".join(MCP_MODES))
    return value


def agent_dir(agent):
    """Resolve an agent id to its folder, refusing anything else."""
    if agent not in agent_folders():
        raise ValueError(f"unknown agent {agent!r}")
    return os.path.join(AGENTS_DIR, agent)


def _resolve_under(root, *parts):
    """Resolve an existing or prospective path beneath a canonical root."""
    root = os.path.realpath(root)
    target = os.path.realpath(os.path.join(root, *parts))
    try:
        contained = os.path.commonpath((root, target)) == root
    except ValueError:
        contained = False
    if not contained:
        raise ValueError("requested path is outside its allowed root")
    return target


def available_domains():
    root = os.path.join(PROJECT, "domains")
    found = []
    if not os.path.isdir(root):
        return found
    for name in sorted(os.listdir(root)):
        if name.startswith("_") or not os.path.isdir(os.path.join(root, name)):
            continue
        try:
            domain = load_domain(name)
        except (ImportError, TypeError, ValueError):
            continue
        found.append({
            "name": domain.name,
            "version": domain.version,
            "presets": list(domain.presets),
        })
    return found


def installed_tags():
    try:
        r = requests.get(f"{OLLAMA_URL}/api/tags", timeout=3)
        r.raise_for_status()
        return {m["name"]: m.get("size", 0) for m in r.json().get("models", [])}
    except Exception:
        return None  # None = server unreachable, {} = up with no models


def tag_installed(tag, tags):
    """Ollama treats llama3.1:8b and llama3.1:latest as different tags even when
    they share blobs, so only an exact match counts — except that a bare name
    means :latest."""
    if not tags:
        return False
    return tag in tags or (":" not in tag and f"{tag}:latest" in tags)


def agent_list():
    tags = installed_tags()
    out = []
    for name in agent_folders():
        cfg = read_config(name)
        folder = os.path.join(AGENTS_DIR, name)
        domain = load_domain(cfg.get("domain") or "office_demo")
        paths = agent_runtime_paths(folder, domain)
        files_dir = str(paths.artifacts)
        logs_dir = str(paths.logs)
        mem_path = str(paths.memory)
        validate_regular_tree_under(folder, paths.root, must_exist=False)
        file_count = 0
        if os.path.isdir(files_dir):
            file_count = len(regular_entries_under(files_dir, limit=5_000))
        run_count = 0
        if os.path.isdir(logs_dir):
            run_count = len(
                regular_entries_under(logs_dir, prefix="run_", suffix=".json")
            )
        memory_count = 0
        if os.path.isfile(mem_path):
            trusted_directory_under(folder, os.path.dirname(mem_path))
            safe_memory = regular_path_under(
                os.path.dirname(mem_path), os.path.basename(mem_path),
                maximum_bytes=MAX_LOG_FILE_BYTES,
            )
            with open(safe_memory, encoding="utf-8") as stream:
                memory_count = sum(1 for _ in stream)
        speed, blurb = SPEED_HINT.get(name, ("", ""))
        out.append({
            "id": name,
            "name": cfg.get("name", name),
            "model": cfg["model"],
            "note": cfg.get("note", ""),
            "domain": domain.name,
            "domain_version": domain.version,
            "presets": list(domain.presets),
            "speed": speed,
            "blurb": blurb,
            "catalog": catalog_for(cfg["model"]),
            "installed": tag_installed(cfg["model"], tags),
            # What this agent connects to when a run does not say otherwise.
            # Without it the options panel reads "Real accounts: none" while a
            # config quietly enables a live mailbox for every run.
            "mcp_default": list((cfg.get("mcp") or {}).get("enable") or []),
            "mcp_default_mode": (cfg.get("mcp") or {}).get("mode") or "draft",
            "files": file_count,
            "runs": run_count,
            "memories": memory_count,
        })
    presets = out[0]["presets"] if out else []
    # Models the catalog knows about that are not installed. The rail offers
    # the pull command for them, so a machine with one model is not a dead end.
    have = {a["model"] for a in out}
    offered = [dict(value, tag=key) for key, value in CATALOG.items()
               if key not in have and not tag_installed(key, tags)]
    offered.sort(key=lambda m: m["tag"])
    return {"agents": out, "domains": available_domains(),
            "ollama": tags is not None, "presets": presets,
            "available": offered, "installed_models": sorted(tags or {}),
            "project": PROJECT}


# -------------------------------------------------------------- workspace ----

def _agent_domain(agent, domain_name=None):
    config = read_config(agent)
    return load_domain(domain_name or config.get("domain") or "office_demo")


def workspace(agent, domain_name=None):
    """The agent's folder as the browser shows it — same shape the runner emits
    during a run, so the panel renders identically live and at rest."""
    folder = agent_dir(agent)
    domain = _agent_domain(agent, domain_name)
    paths = agent_runtime_paths(folder, domain)
    validate_regular_tree_under(folder, paths.root, must_exist=False)
    state = domain.inspect(paths.workspace, paths.memory)

    logs = []
    if os.path.isdir(paths.logs):
        for mtime_ns, name, _size in regular_entries_under(
            paths.logs, prefix="run_", suffix=".json", limit=25
        ):
            logs.append({"name": name, "mtime": mtime_ns / 1_000_000_000})
    state["logs"] = logs
    state["folder"] = str(paths.root)
    return state


def workspace_file(agent, name, domain_name=None):
    folder = agent_dir(agent)
    domain = _agent_domain(agent, domain_name)
    files_dir = str(agent_runtime_paths(folder, domain).artifacts)
    trusted_directory_under(folder, files_dir)
    return regular_path_under(files_dir, name, maximum_bytes=MAX_ARTIFACT_BYTES)


def _check_office_archive(path):
    """Reject oversized/degenerate Office archives before library expansion."""
    with zipfile.ZipFile(path) as archive:
        members = archive.infolist()
        if len(members) > MAX_ARCHIVE_MEMBERS:
            raise RequestError(413, "preview archive has too many members")
        expanded = sum(member.file_size for member in members)
        if expanded > MAX_ARCHIVE_EXPANDED_BYTES:
            raise RequestError(413, "preview archive expands beyond its limit")


def preview(agent, name, domain_name=None):
    """Render a generated file in the browser instead of making the user open
    PowerPoint — the whole point is to see what the agent produced."""
    path = workspace_file(agent, name, domain_name)
    if os.path.getsize(path) > MAX_PREVIEW_BYTES:
        raise RequestError(413, "file is too large to preview")
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pptx":
        _check_office_archive(path)
        from pptx import Presentation
        slides = []
        for slide in list(Presentation(path).slides)[:100]:
            title, body = "", []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                lines = [p.text for p in shape.text_frame.paragraphs if p.text.strip()]
                if shape == slide.shapes.title:
                    title = " ".join(lines)
                else:
                    body += lines[:200]
            slides.append({"title": title[:4_096], "bullets": body[:500]})
        return {"kind": "pptx", "name": os.path.basename(path), "slides": slides}
    if ext == ".xlsx":
        _check_office_archive(path)
        from openpyxl import load_workbook
        sheets = []
        book = load_workbook(path, data_only=False, read_only=True)
        for ws in book.worksheets[:20]:
            rows = []
            for row in ws.iter_rows(max_row=200, max_col=50, values_only=True):
                rows.append([
                    "" if cell is None else str(cell)[:4_096] for cell in row
                ])
            sheets.append({"sheet": ws.title[:256], "rows": rows})
        book.close()
        return {"kind": "xlsx", "name": os.path.basename(path), "sheets": sheets}
    if ext not in (".txt", ".json", ".md", ".csv"):
        return {"kind": "binary", "name": os.path.basename(path),
                "size": os.path.getsize(path)}
    with open(path, "rb") as f:
        blob = f.read(20000)
    if b"\x00" in blob[:2000]:
        return {"kind": "binary", "name": os.path.basename(path),
                "size": os.path.getsize(path)}
    return {"kind": "text", "name": os.path.basename(path),
            "text": blob.decode("utf-8", errors="replace")}


def reset_agent(agent, what, domain_name=None):
    folder = agent_dir(agent)
    domain = _agent_domain(agent, domain_name)
    paths = agent_runtime_paths(folder, domain)
    done = []
    targets = {
        "world": str(paths.workspace / "state.json"),
        "memory": str(paths.memory),
    }
    for key, path in targets.items():
        if key in what and os.path.isfile(path):
            os.remove(path)
            done.append(key)
    for key, path in (
        ("files", str(paths.artifacts)),
        ("logs", str(paths.logs)),
    ):
        if key in what and os.path.isdir(path):
            reset_directory(folder, path)
            done.append(key)
    return done


def reveal(path):
    if sys.platform == "darwin":
        subprocess.Popen(["open", path])
    elif sys.platform.startswith("win"):
        os.startfile(path)  # noqa: S606 - local dev convenience
    else:
        subprocess.Popen(["xdg-open", path])


# -------------------------------------------------------------------- runs ----

def thread_reply(end, status):
    """What the agent says back in the conversation.

    Its done() summary when there is one. A small model often never calls
    done(), and the placeholder that used to stand in for that said only that
    there was no summary, so a conversation where real work happened read as a
    row of shrugs. Report the steps instead: the run drafted a mail, and the
    transcript should say so."""
    if end and end.get("summary"):
        return end["summary"]
    if status == "stopped":
        return "(stopped before finishing)"
    done = []
    for action in (end or {}).get("actions") or []:
        name = action.get("tool")
        if not action.get("ok") or not name:
            continue
        if done and done[-1][0] == name:
            done[-1][1] += 1
        else:
            done.append([name, 1])
    if not done:
        return "(no summary, and no step completed)"
    steps = ", ".join(n if c == 1 else f"{n} x{c}" for n, c in done)
    return f"(no summary) Steps completed: {steps}."


class Run:
    """One agent subprocess, its event log, and everyone watching it."""

    def __init__(self, rid, agent, task, process_tree, options):
        self.id = rid
        self.agent = agent
        self.task = task
        self.process_tree = process_tree
        self.proc = process_tree.proc
        self.options = options
        self.started = time.time()
        self.events = EventJournal()
        self.confirmations = ConfirmationLedger(rid)
        self.status = "running"
        self.lock = threading.Lock()

    def add(self, event):
        return self.events.add(event)

    def subscribe(self, after=-1):
        """Register a watcher and hand back everything it has not seen. `after`
        comes from Last-Event-ID, so a reconnect resumes instead of replaying."""
        return self.events.subscribe(after)

    def unsubscribe(self, q):
        self.events.unsubscribe(q)

    def stop(self):
        with self.lock:
            if self.proc.poll() is not None:
                return
            self.status = "stopped"
            self.confirmations.clear()
            try:
                self.proc.stdin.close()
            except (AttributeError, OSError):
                pass
        self.process_tree.terminate()

    def register_confirmation(self, event):
        if event.get("run_id") != self.id:
            raise ValueError("confirmation belongs to another run")
        self.confirmations.register(
            event.get("confirmation_id"), event.get("nonce")
        )

    def decide(self, confirmation_id, nonce, decision):
        message = self.confirmations.decide(
            self.id, confirmation_id, nonce, decision
        )
        if self.proc.poll() is not None:
            raise RequestError(409, "run is no longer active")
        try:
            self.proc.stdin.write(json.dumps(message) + "\n")
            self.proc.stdin.flush()
        except (BrokenPipeError, OSError):
            raise RequestError(409, "run confirmation channel is closed")


class Runs:
    def __init__(self):
        self.current = None
        self.lock = threading.Lock()

    def start(self, agent, task, options):
        with self.lock:
            if self.current and self.current.proc.poll() is None:
                raise RuntimeError(f"{self.current.agent} is already running — "
                                   "stop it first (one agent run at a time).")
            run_id = new_capability()[:22]
            cmd = [sys.executable, "-u", "-m", "webui.runner",
                   "--run-id", run_id, "--agent", agent, "--task", task]
            if options.get("domain"):
                cmd += ["--domain", options["domain"]]
            if options.get("tiers"):
                cmd.append("--tiers")
            if options.get("small"):
                cmd += ["--small", options["small"]]
            if options.get("deep"):
                cmd += ["--deep", options["deep"]]
            if options.get("max_calls") is not None:
                cmd += ["--max-calls", str(int(options["max_calls"]))]
            if options.get("thread"):
                cmd += ["--thread", options["thread"]]
            if options.get("model"):
                cmd += ["--model", options["model"]]
            if options.get("mcp"):
                cmd += ["--mcp", ",".join(options["mcp"])]
                if options.get("mcp_mode"):
                    cmd += ["--mcp-mode", options["mcp_mode"]]
                if options.get("keep_office_tools"):
                    cmd.append("--keep-office-tools")
            env = dict(os.environ, PYTHONUNBUFFERED="1", PYTHONIOENCODING="utf-8")
            process_tree = ProcessTree.start(
                cmd, cwd=PROJECT, env=env, text=True,
                encoding="utf-8", errors="replace", bufsize=1,
                stdin=subprocess.PIPE, stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
            )
            run = Run(run_id, agent, task, process_tree, options)
            self.current = run
        threading.Thread(target=self._pump, args=(run,), daemon=True).start()
        return run

    def _pump(self, run):
        stderr = deque(maxlen=MAX_STDERR_LINES)

        def drain_stderr():
            for line in run.proc.stderr:
                stderr.append(line[-MAX_STDERR_LINE:])

        stderr_thread = threading.Thread(target=drain_stderr, daemon=True)
        stderr_thread.start()
        for line in run.proc.stdout:
            line = line.strip()
            if not line:
                continue
            try:
                event = json.loads(line)
                if not isinstance(event, dict):
                    raise ValueError("runner event is not an object")
                if event.get("t") == "confirmation":
                    run.register_confirmation(event)
                run.add(event)
            except ValueError:
                run.add({"t": "stdout", "text": line[:MAX_STDERR_LINE]})
        code = run.proc.wait()
        stderr_thread.join(timeout=0.2)
        run.process_tree.close()
        run.confirmations.clear()
        if run.status == "running":
            run.status = "finished" if code == 0 else "failed"
        if code not in (0, -15) and run.status != "stopped":
            tail = "".join(stderr)[-1500:].strip()
            if tail:
                safe_tail = redact(tail)
                sys.stderr.write(
                    f"Agent Lab runner {run.id} exited with code {code}: "
                    f"{safe_tail}\n"
                )
                sys.stderr.flush()
            run.add({"t": "error", "message": f"the run exited with code {code}"})
        # The agent's reply to the conversation is its done() summary. Recorded
        # here rather than in the runner so a crashed or stopped run still
        # leaves an honest turn in the thread instead of a silent gap.
        thread_id = run.options.get("thread")
        if thread_id:
            end = next((event for _, event in reversed(run.events.snapshot())
                        if event.get("t") == "end"), None)
            reply = thread_reply(end, run.status)
            chat.append(agent_dir(run.agent), thread_id, "assistant", reply,
                        run=run.id)
        run.add({"t": "closed", "status": run.status, "code": code})

    def with_idle(self, operation):
        """Serialize reset with process creation and reject a live run."""
        with self.lock:
            if self.current and self.current.proc.poll() is None:
                raise RuntimeError("cannot reset state while a run is active")
            return operation()


RUNS = Runs()


# ------------------------------------------------------------------ server ----

class Handler(http.server.BaseHTTPRequestHandler):
    protocol_version = "HTTP/1.1"
    server_version = "AgentLab"

    def log_message(self, *args):
        pass  # the console belongs to the run banner, not to request noise

    # ---- helpers ----
    def _security_headers(self):
        self.send_header("Cache-Control", "no-store")
        self.send_header("Content-Security-Policy", "default-src 'self'; "
                         "connect-src 'self'; img-src 'self' data:; "
                         "style-src 'self'; script-src 'self'; object-src 'none'; "
                         "base-uri 'none'; frame-ancestors 'none'")
        self.send_header("Cross-Origin-Opener-Policy", "same-origin")
        self.send_header("Cross-Origin-Resource-Policy", "same-origin")
        self.send_header("Referrer-Policy", "no-referrer")
        self.send_header("X-Content-Type-Options", "nosniff")
        self.send_header("X-Frame-Options", "DENY")

    def send_json(self, obj, status=200):
        body = json.dumps(obj, ensure_ascii=False, default=str).encode("utf-8")
        self.send_response(status)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.send_header("Content-Length", str(len(body)))
        self._security_headers()
        self.end_headers()
        self.wfile.write(body)

    def send_bytes(self, blob, ctype, extra=(), status=200):
        self.send_response(status)
        self.send_header("Content-Type", ctype)
        self.send_header("Content-Length", str(len(blob)))
        self._security_headers()
        for k, v in extra:
            self.send_header(k, v)
        self.end_headers()
        self.wfile.write(blob)

    def query(self):
        parts = urllib.parse.urlparse(self.path)
        pairs = urllib.parse.parse_qsl(
            parts.query, keep_blank_values=True, strict_parsing=True
        ) if parts.query else []
        values = {}
        for key, value in pairs:
            if key in values:
                raise RequestError(400, "duplicate query parameter")
            values[key] = value
        return parts.path, values

    def exact_query(self, query, *, required=(), optional=()):
        return exact_object(query, required=required, optional=optional)

    def authorize(self, path, *, mutation=False):
        validate_host(self.headers, self.server.expected_host)
        if path.startswith("/api/"):
            validate_capability(self.headers, self.server.capability)
        if mutation:
            validate_mutation_origin(self.headers, self.server.origin)

    def handle_error(self, exc):
        if isinstance(exc, RequestError):
            if exc.status == 401:
                self.send_response(401)
                self.send_header("WWW-Authenticate", 'Bearer realm="Agent Lab"')
                body = json.dumps({"error": str(exc)}).encode("utf-8")
                self.send_header("Content-Type", "application/json; charset=utf-8")
                self.send_header("Content-Length", str(len(body)))
                self._security_headers()
                self.end_headers()
                self.wfile.write(body)
                return
            return self.send_json({"error": str(exc)}, exc.status)
        if isinstance(exc, RuntimeError):
            return self.send_json({"error": str(exc)}, 409)
        if isinstance(exc, ValueError):
            return self.send_json({"error": str(exc)}, 400)
        if isinstance(exc, FileNotFoundError):
            return self.send_json({"error": "file not found"}, 404)
        return self.send_json({"error": "internal Agent Lab error"}, 500)

    # ---- GET ----
    def do_GET(self):
        try:
            path, q = self.query()
            self.authorize(path)
            if path in ("/", "/index.html"):
                self.exact_query(q)
                return self.static_file("index.html")
            if path.startswith("/static/"):
                self.exact_query(q)
                return self.static_file(path[len("/static/"):])
            # Both must be served from the root, not /static/: a service worker
            # may only control pages at or below its own path, and the manifest
            # is resolved relative to the document.
            if path == "/sw.js":
                self.exact_query(q)
                return self.static_file("sw.js")
            if path == "/manifest.webmanifest":
                self.exact_query(q)
                return self.static_file("manifest.webmanifest")
            if path == "/api/agents":
                self.exact_query(q)
                return self.send_json(agent_list())
            if path == "/api/workspace":
                q = self.exact_query(q, required=("agent",), optional=("domain",))
                return self.send_json(
                    workspace(q.get("agent", ""), q.get("domain"))
                )
            if path == "/api/preview":
                q = self.exact_query(
                    q, required=("agent", "name"), optional=("domain",)
                )
                return self.send_json(
                    preview(
                        q.get("agent", ""),
                        q.get("name", ""),
                        q.get("domain"),
                    )
                )
            if path == "/api/download":
                q = self.exact_query(
                    q, required=("agent", "name"), optional=("domain",)
                )
                fpath = workspace_file(
                    q.get("agent", ""),
                    q.get("name", ""),
                    q.get("domain"),
                )
                ctype = mimetypes.guess_type(fpath)[0] or "application/octet-stream"
                with open(fpath, "rb") as f:
                    blob = f.read()
                return self.send_bytes(blob, ctype, extra=[
                    ("Content-Disposition",
                     f'attachment; filename="{os.path.basename(fpath)}"')])
            if path == "/api/log":
                q = self.exact_query(
                    q, required=("agent", "name"), optional=("domain",)
                )
                folder = agent_dir(q.get("agent", ""))
                domain = _agent_domain(
                    q.get("agent", ""), q.get("domain")
                )
                paths = agent_runtime_paths(folder, domain)
                trusted_directory_under(folder, paths.logs)
                name = portable_leaf(q.get("name", ""))
                if not name.startswith("run_") or not name.endswith(".json"):
                    raise RequestError(400, "invalid run-log name")
                log_path = regular_path_under(
                    paths.logs, name, maximum_bytes=MAX_LOG_FILE_BYTES
                )
                with open(log_path, encoding="utf-8") as f:
                    return self.send_json(json.load(f))
            if path == "/api/threads":
                return self.send_json(
                    {"threads": chat.threads(agent_dir(q.get("agent", "")))})
            if path == "/api/thread":
                folder = agent_dir(q.get("agent", ""))
                return self.send_json(
                    {"id": q.get("id", ""),
                     "messages": chat.messages(folder, q.get("id", ""))})
            if path == "/api/mcp":
                # The registry, for the run panel's account picker. Setup notes
                # come along so the UI can say what a server needs before it
                # works, rather than failing at connect time.
                return self.send_json([
                    {"name": name, "summary": summary,
                     "setup": mcp_config.setup_notes(name)}
                    for name, summary in mcp_config.available()])
            if path == "/api/status":
                self.exact_query(q)
                run = self.server.runs.current
                return self.send_json({"run": run.id if run else None,
                                       "agent": run.agent if run else None,
                                       "status": run.status if run else "idle"})
            if path == "/api/events":
                q = self.exact_query(q, required=("run",))
                return self.stream_events(q)
        except Exception as exc:
            return self.handle_error(exc)
        self.send_json({"error": "not found"}, 404)

    # ---- POST ----
    def do_POST(self):
        try:
            path, query = self.query()
            self.exact_query(query)
            self.authorize(path, mutation=True)
            body = read_json_object(self)
            if path == "/api/thread/new":
                body = exact_object(body, required=("agent",),
                                    optional=("task",))
                folder = agent_dir(require_string(body["agent"], "agent",
                                                  maximum=128))
                task = require_optional_string(body.get("task"), "task",
                                               maximum=8_192) or ""
                return self.send_json({"id": chat.create(folder, task)})
            if path == "/api/thread/delete":
                body = exact_object(body, required=("agent", "id"))
                folder = agent_dir(require_string(body["agent"], "agent",
                                                  maximum=128))
                thread_id = require_string(body["id"], "id", maximum=128)
                return self.send_json(
                    {"deleted": chat.delete(folder, thread_id)})
            if path == "/api/run":
                reject_removed_run_fields(body)
                body = exact_object(
                    body,
                    required=("agent", "domain", "task", "tiers", "max_calls"),
                    optional=("small", "deep", "mcp", "mcp_mode",
                              "keep_office_tools", "thread", "model"),
                )
                agent = require_string(body["agent"], "agent", maximum=128)
                task = require_string(body["task"], "task").strip()
                domain = require_optional_string(body["domain"], "domain")
                tiers = require_bool(body["tiers"], "tiers")
                max_calls = body["max_calls"]
                if max_calls is not None:
                    max_calls = require_int(
                        max_calls, "max_calls", minimum=2, maximum=80
                    )
                agent_dir(agent)  # validates
                if not task:
                    raise RequestError(400, "give the agent a task first")
                options = {
                    "thread": require_optional_string(body.get("thread"),
                                                      "thread", maximum=128),
                    "model": require_optional_string(body.get("model"),
                                                     "model", maximum=200),
                    "mcp": require_mcp_names(body.get("mcp")),
                    "mcp_mode": require_mcp_mode(body.get("mcp_mode")),
                    "keep_office_tools": require_bool(
                        body.get("keep_office_tools") or False,
                        "keep_office_tools"),
                    "domain": domain,
                    "tiers": tiers,
                    "small": require_optional_string(body.get("small"), "small"),
                    "deep": require_optional_string(body.get("deep"), "deep"),
                    "max_calls": max_calls,
                }
                if options["thread"]:
                    chat.append(agent_dir(agent), options["thread"],
                                "user", task)
                run = self.server.runs.start(agent, task, options)
                return self.send_json({"run": run.id, "agent": agent})
            if path == "/api/stop":
                body = exact_object(body, required=("run_id",))
                run_id = require_string(body["run_id"], "run_id", maximum=128)
                run = self.server.runs.current
                if not run or run.id != run_id:
                    raise RequestError(409, "run is stale or unknown")
                run.stop()
                return self.send_json({"ok": True})
            if path == "/api/reset":
                body = exact_object(
                    body, required=("agent", "domain", "what")
                )
                agent = require_string(body["agent"], "agent", maximum=128)
                domain = require_optional_string(body["domain"], "domain")
                values = body["what"]
                if (
                    not isinstance(values, list)
                    or not values
                    or any(not isinstance(item, str) for item in values)
                    or len(values) != len(set(values))
                ):
                    raise RequestError(400, "what must be a nonempty unique string list")
                what = set(values)
                if not what <= {"world", "memory", "files", "logs"}:
                    raise RequestError(400, "what contains an unsupported reset target")
                return self.send_json({
                    "cleared": self.server.runs.with_idle(
                        lambda: reset_agent(agent, what, domain)
                    )
                })
            if path == "/api/reveal":
                body = exact_object(body, required=("agent", "domain"))
                agent = require_string(body["agent"], "agent", maximum=128)
                folder = agent_dir(agent)
                domain = _agent_domain(
                    agent, require_optional_string(body["domain"], "domain")
                )
                base = str(agent_runtime_paths(folder, domain).root)
                trusted_directory_under(folder, base)
                reveal(base)
                return self.send_json({"ok": True})
            if path == "/api/confirm":
                body = exact_object(
                    body,
                    required=("run_id", "confirmation_id", "nonce", "decision"),
                )
                run_id = require_string(body["run_id"], "run_id", maximum=128)
                run = self.server.runs.current
                if not run or run.id != run_id:
                    raise RequestError(409, "confirmation belongs to a stale run")
                run.decide(
                    require_string(
                        body["confirmation_id"], "confirmation_id", maximum=128
                    ),
                    require_string(body["nonce"], "nonce", minimum=32, maximum=256),
                    require_bool(body["decision"], "decision"),
                )
                return self.send_json({"ok": True})
            if path == "/api/pull":
                body = exact_object(body, required=("model",))
                return self.stream_pull(
                    require_string(body["model"], "model", maximum=200)
                )
        except Exception as exc:
            return self.handle_error(exc)
        self.send_json({"error": "not found"}, 404)

    # ---- static ----
    def static_file(self, rel):
        try:
            path = _resolve_under(STATIC, rel)
        except ValueError:
            return self.send_json({"error": "not found"}, 404)
        if not os.path.isfile(path):
            return self.send_json({"error": "not found"}, 404)
        with open(path, "rb") as f:
            blob = f.read()
        ctype = mimetypes.guess_type(path)[0] or "text/plain"
        self.send_bytes(blob, f"{ctype}; charset=utf-8" if "text" in ctype
                        or "javascript" in ctype else ctype)

    # ---- SSE ----
    def open_stream(self):
        self.send_response(200)
        self.send_header("Content-Type", "text/event-stream; charset=utf-8")
        self._security_headers()
        self.send_header("Connection", "close")
        self.end_headers()

    def push(self, obj, index=None):
        prefix = f"id: {index}\n" if index is not None else ""
        self.wfile.write(f"{prefix}data: {json.dumps(obj, ensure_ascii=False, default=str)}\n\n"
                         .encode("utf-8"))
        self.wfile.flush()

    def stream_events(self, q):
        run = self.server.runs.current
        want = q.get("run")
        if not run or (want and str(run.id) != str(want)):
            self.open_stream()
            self.push({"t": "closed", "status": "gone"})
            return
        try:
            after = int(self.headers.get("Last-Event-ID"))
        except (TypeError, ValueError):
            after = -1
        sub, backlog = run.subscribe(after)
        self.open_stream()
        self.close_connection = True
        try:
            for index, event in backlog:
                self.push(event, index)
            while True:
                try:
                    index, event = sub.get(timeout=10)
                except queue.Empty:
                    self.wfile.write(b": keepalive\n\n")
                    self.wfile.flush()
                    continue
                self.push(event, index)
                if event.get("t") == "closed":
                    break
        except (BrokenPipeError, ConnectionResetError):
            pass
        finally:
            run.unsubscribe(sub)

    def stream_pull(self, model):
        """Download a model from the picker, so 'run the 14B' never means
        leaving the page for a terminal."""
        self.open_stream()
        self.close_connection = True
        if not model:
            return self.push({"t": "error", "message": "no model given"})
        try:
            with requests.post(f"{OLLAMA_URL}/api/pull", json={"model": model},
                               stream=True, timeout=None) as r:
                r.raise_for_status()
                for line in r.iter_lines(decode_unicode=True):
                    if not line:
                        continue
                    msg = json.loads(line)
                    self.push({"t": "pull", "status": msg.get("status", ""),
                               "completed": msg.get("completed", 0),
                               "total": msg.get("total", 0),
                               "error": msg.get("error")})
            self.push({"t": "closed", "status": "done"})
        except (BrokenPipeError, ConnectionResetError):
            pass
        except Exception as exc:
            detail = redact(f"{type(exc).__name__}: {exc}")
            sys.stderr.write(f"Agent Lab model pull failed: {detail}\n")
            sys.stderr.flush()
            try:
                self.push({"t": "error", "message": "model pull failed"})
            except OSError:
                pass


class Server(http.server.ThreadingHTTPServer):
    daemon_threads = True
    allow_reuse_address = True

    def __init__(self, server_address, handler_class, *, capability=None, runs=None):
        super().__init__(server_address, handler_class)
        host, port = self.server_address[:2]
        self.expected_host = f"{host}:{port}"
        self.origin = f"http://{self.expected_host}"
        self.capability = capability or new_capability()
        self.runs = runs or Runs()


def bind_server(start=DEFAULT_PORT):
    """Bind directly, avoiding the check-then-bind race of port probing."""
    if start == 0:
        return Server(("127.0.0.1", 0), Handler)
    last_error = None
    for port in range(start, start + 20):
        try:
            return Server(("127.0.0.1", port), Handler)
        except OSError as exc:
            last_error = exc
    raise last_error


def serve_until_stopped(server):
    """Serve and always close the active process tree and listening socket."""
    interrupted = False
    serve_error = None
    serve_traceback = None
    try:
        server.serve_forever()
    except KeyboardInterrupt:
        interrupted = True
    except BaseException as exc:  # Preserve the serving failure after cleanup.
        serve_error = exc
        serve_traceback = exc.__traceback__
    cleanup_error = None
    try:
        current = server.runs.current
        if current and current.proc.poll() is None:
            current.stop()
    except Exception as exc:
        cleanup_error = exc
        detail = redact(f"{type(exc).__name__}: {exc}")
        sys.stderr.write(f"Agent Lab run cleanup failed: {detail}\n")
        sys.stderr.flush()
    try:
        server.server_close()
    except Exception as exc:
        if cleanup_error is None:
            cleanup_error = exc
        detail = redact(f"{type(exc).__name__}: {exc}")
        sys.stderr.write(f"Agent Lab socket cleanup failed: {detail}\n")
        sys.stderr.flush()
    if serve_error is not None:
        raise serve_error.with_traceback(serve_traceback)
    if cleanup_error is not None:
        raise RuntimeError("Agent Lab cleanup did not complete") from cleanup_error
    return interrupted


def main():
    sys.path.insert(0, PROJECT)
    requested_port = int(os.environ.get("AGENT_LAB_PORT", DEFAULT_PORT))
    server = bind_server(requested_port)
    url = server.origin + "/#capability=" + server.capability
    tags = installed_tags()
    print(f"\n  Agent Lab  →  {url}")
    print(f"  project    {PROJECT}")
    if tags is None:
        print("  ollama     NOT RUNNING — start it, then reload the page")
    else:
        print(f"  ollama     up, {len(tags)} model(s) installed")
    print("  agents     " + ", ".join(agent_folders()))
    print("\n  Ctrl-C to stop.\n")
    if os.environ.get("AGENT_LAB_NO_BROWSER") != "1":
        threading.Timer(0.6, lambda: webbrowser.open(url)).start()
    if serve_until_stopped(server):
        print("\n  stopped.\n")


if __name__ == "__main__":
    main()
