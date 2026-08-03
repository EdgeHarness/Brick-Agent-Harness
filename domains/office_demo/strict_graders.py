"""S5 strict graders for the fixed synthetic office scenarios."""
from collections import Counter
from io import BytesIO
import re

from openpyxl import load_workbook
from pptx import Presentation

from harness.grading import GraderSpec, GradingError


GRADER_VERSION = "1.0.0"
CHECKS = (
    ("required_outcome", "the complete requested outcome is correct"),
    ("source_observed", "required source state was inspected"),
    ("exact_business_effects", "exactly the requested business effects occurred"),
    ("exact_artifacts", "the artifact set has exactly the requested names"),
    ("no_unrequested_state", "no unrequested authoritative state changed"),
)
BUSINESS_TOOLS = frozenset(
    {
        "send_email",
        "add_event",
        "send_message",
        "set_reminder",
        "create_presentation",
        "create_spreadsheet",
    }
)
NONBUSINESS_TOOLS = frozenset(
    {
        "list_emails",
        "read_email",
        "list_events",
        "read_spreadsheet",
        "think",
        "save_memory",
        "recall_memories",
        "done",
    }
)


def _inputs(evidence):
    if evidence.domain != "office_demo":
        raise GradingError("office grader received another domain")
    state = evidence.state
    required = {
        "emails",
        "events",
        "sent_emails",
        "messages",
        "reminders",
    }
    if not isinstance(state, dict) or set(state) != required:
        raise GradingError("office grading state has the wrong schema")
    actions = evidence.actions
    if not isinstance(actions, list):
        raise GradingError("actions must be a list")
    for action in actions:
        if (
            not isinstance(action, dict)
            or not isinstance(action.get("tool"), str)
            or type(action.get("ok")) is not bool
            or not isinstance(action.get("args"), dict)
        ):
            raise GradingError("action evidence has the wrong schema")
    return state, actions, dict(evidence.artifact_map())


def _successful(actions, tool=None):
    return [
        item
        for item in actions
        if item["ok"] and (tool is None or item["tool"] == tool)
    ]


def _read(actions, tool, predicate=lambda _args: True):
    return any(
        item["tool"] == tool and item["ok"] and predicate(item["args"])
        for item in actions
    )


def _effects(actions, expected):
    if any(
        item["ok"]
        and item["tool"] not in BUSINESS_TOOLS
        and item["tool"] not in NONBUSINESS_TOOLS
        for item in actions
    ):
        return False
    actual = Counter(
        item["tool"]
        for item in actions
        if item["ok"] and item["tool"] in BUSINESS_TOOLS
    )
    return actual == Counter(expected)


def _state_exact(state, *, events=7, emails=0, messages=0, reminders=0):
    return (
        len(state["emails"]) == 10
        and len(state["events"]) == events
        and len(state["sent_emails"]) == emails
        and len(state["messages"]) == messages
        and len(state["reminders"]) == reminders
    )


def _base(outcome, source, effects, artifacts, state):
    return {
        "required_outcome": bool(outcome),
        "source_observed": bool(source),
        "exact_business_effects": bool(effects),
        "exact_artifacts": bool(artifacts),
        "no_unrequested_state": bool(state),
    }


def _slides(payload):
    try:
        deck = Presentation(BytesIO(payload))
    except Exception as exc:
        raise GradingError(f"PowerPoint artifact is corrupt: {exc}") from exc
    result = []
    for slide in deck.slides:
        title = ""
        bullets = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            texts = [
                paragraph.text.strip()
                for paragraph in shape.text_frame.paragraphs
                if paragraph.text.strip()
            ]
            if shape == slide.shapes.title:
                title = " ".join(texts)
            else:
                bullets.extend(texts)
        result.append((title, bullets))
    return result


def _rows(payload):
    try:
        book = load_workbook(BytesIO(payload), data_only=False, read_only=True)
        sheet = book.worksheets[0]
        rows = [list(row) for row in sheet.iter_rows(values_only=True)]
        book.close()
    except Exception as exc:
        raise GradingError(f"spreadsheet artifact is corrupt: {exc}") from exc
    while rows and not any(value is not None for value in rows[-1]):
        rows.pop()
    return rows


def _text(value):
    return re.sub(r"\s+", " ", str(value or "").strip().casefold())


def _number(value, rows=None):
    if isinstance(value, bool):
        return None
    if isinstance(value, (int, float)):
        return float(value)
    raw = str(value or "").strip().replace("$", "").replace(",", "")
    match = re.fullmatch(r"=SUM\(([A-Z]+)(\d+):([A-Z]+)(\d+)\)", raw, re.I)
    if match and rows is not None:
        if match.group(1).casefold() != match.group(3).casefold():
            return None
        column = 0
        for character in match.group(1).upper():
            column = column * 26 + ord(character) - ord("A") + 1
        column -= 1
        start, end = int(match.group(2)) - 1, int(match.group(4)) - 1
        values = [
            _number(rows[index][column])
            for index in range(start, end + 1)
            if index < len(rows) and column < len(rows[index])
        ]
        if all(value is not None for value in values):
            return sum(values)
    try:
        return float(raw)
    except ValueError:
        return None


def _artifact_case(evidence, name, expected_effect, source, validate, state_ok):
    state, actions, artifacts = _inputs(evidence)
    payload = artifacts.get(name)
    outcome = False if payload is None else validate(payload)
    return _base(
        outcome,
        source(actions),
        _effects(actions, [expected_effect]),
        set(artifacts) == {name},
        state_ok(state),
    )


def pptx_basic(evidence):
    def validate(payload):
        slides = _slides(payload)
        titles = [
            "q3 business review",
            "agenda",
            "sales",
            "marketing",
            "next steps",
        ]
        return (
            len(slides) == 5
            and [_text(item[0]) for item in slides] == titles
            and all(len(slides[index][1]) >= 3 for index in range(1, 5))
        )

    return _artifact_case(
        evidence,
        "q3_review.pptx",
        "create_presentation",
        lambda _actions: True,
        validate,
        _state_exact,
    )


def pptx_from_email(evidence):
    expected = {"west": "1240000", "east": "845000", "online": "610000"}

    def validate(payload):
        slides = _slides(payload)
        if len(slides) != 4:
            return False
        regions = {}
        for title, bullets in slides[1:]:
            region = _text(title)
            if region not in expected or region in regions:
                return False
            regions[region] = re.sub(r"[,\s$]", "", " ".join(bullets).casefold())
        return set(regions) == set(expected) and all(
            number in regions[region] for region, number in expected.items()
        )

    return _artifact_case(
        evidence,
        "sales_summary.pptx",
        "create_presentation",
        lambda actions: _read(actions, "read_email", lambda args: args.get("id") == "e2"),
        validate,
        _state_exact,
    )


def xlsx_basic(evidence):
    wanted = {
        "laptops": 3200.0,
        "software licenses": 1150.0,
        "training": 800.0,
        "travel": 2400.0,
    }

    def validate(payload):
        rows = _rows(payload)
        if len(rows) != 6 or any(len(row) != 2 for row in rows):
            return False
        if [_text(rows[0][0]), _text(rows[0][1])] != ["item", "cost"]:
            return False
        actual = {_text(row[0]): _number(row[1], rows) for row in rows[1:5]}
        return (
            actual == wanted
            and _text(rows[5][0]) == "total"
            and _number(rows[5][1], rows) == 7550.0
        )

    return _artifact_case(
        evidence,
        "budget.xlsx",
        "create_spreadsheet",
        lambda _actions: True,
        validate,
        _state_exact,
    )


def xlsx_from_email(evidence):
    wanted = {
        ("2026-07-12", "cloudhost"): 230.0,
        ("2026-07-14", "officemax"): 87.5,
        ("2026-07-15", "delta"): 412.3,
    }

    def validate(payload):
        rows = _rows(payload)
        if len(rows) != 5 or any(len(row) != 3 for row in rows):
            return False
        if [_text(value) for value in rows[0][:3]] != ["date", "vendor", "amount"]:
            return False
        actual = {
            (_text(row[0]), _text(row[1])): _number(row[2], rows)
            for row in rows[1:4]
        }
        return (
            actual == wanted
            and _text(rows[4][0]) == "total"
            and _number(rows[4][2], rows) == 729.8
        )

    return _artifact_case(
        evidence,
        "expenses.xlsx",
        "create_spreadsheet",
        lambda actions: all(
            _read(actions, "read_email", lambda args, item=item: args.get("id") == item)
            for item in ("e3", "e4", "e5")
        ),
        validate,
        _state_exact,
    )


def email_reply(evidence):
    state, actions, artifacts = _inputs(evidence)
    sent = state["sent_emails"]
    valid = False
    if len(sent) == 1:
        message = sent[0]
        blob = _text(message.get("subject")) + " " + _text(message.get("body"))
        valid = (
            _text(message.get("to")) == "mia@corp.com"
            and ("northwind" in blob or "kickoff" in blob)
            and any(term in blob for term in ("confirm", "attend", "i'll be there", "i will be there", "count me in"))
        )
    return _base(
        valid,
        _read(actions, "read_email", lambda args: args.get("id") == "e7"),
        _effects(actions, ["send_email"]),
        not artifacts,
        _state_exact(state, emails=1),
    )


def _new_events(state):
    return [event for event in state["events"] if event.get("id") not in {f"c{i}" for i in range(1, 8)}]


def cal_add(evidence):
    state, actions, artifacts = _inputs(evidence)
    events = _new_events(state)
    valid = len(events) == 1
    if valid:
        event = events[0]
        valid = (
            _text(event.get("title")) == "design sync"
            and event.get("date") == "2026-07-21"
            and event.get("start") == "14:00"
            and event.get("end") == "15:00"
            and {item.casefold() for item in event.get("attendees", [])}
            == {"alice@corp.com", "bob@corp.com"}
        )
    return _base(valid, True, _effects(actions, ["add_event"]), not artifacts, _state_exact(state, events=8))


def _minutes(value):
    hour, minute = value.split(":")
    return int(hour) * 60 + int(minute)


def cal_freeslot(evidence):
    state, actions, artifacts = _inputs(evidence)
    events = _new_events(state)
    valid = len(events) == 1
    if valid:
        event = events[0]
        start, end = _minutes(event["start"]), _minutes(event["end"])
        occupied = [(540, 660), (720, 780), (900, 960)]
        valid = (
            _text(event.get("title")) == "deep work"
            and event.get("date") == "2026-07-23"
            and end - start == 60
            and 540 <= start < end <= 1020
            and not any(start < other_end and other_start < end for other_start, other_end in occupied)
        )
    source = _read(actions, "list_events", lambda args: args.get("date") in (None, "2026-07-23"))
    return _base(valid, source, _effects(actions, ["add_event"]), not artifacts, _state_exact(state, events=8))


def cal_brief(evidence):
    state, actions, artifacts = _inputs(evidence)
    messages = state["messages"]
    valid = False
    if len(messages) == 1 and "jordan" in _text(messages[0].get("to")):
        body = _text(messages[0].get("text"))
        positions = [body.find(term) for term in ("design review", "1:1 with sam", "marketing sync")]
        valid = -1 < positions[0] < positions[1] < positions[2]
    source = _read(actions, "list_events", lambda args: args.get("date") in (None, "2026-07-22"))
    return _base(valid, source, _effects(actions, ["send_message"]), not artifacts, _state_exact(state, messages=1))


def remind_msg(evidence):
    state, actions, artifacts = _inputs(evidence)
    reminders, messages = state["reminders"], state["messages"]
    valid = (
        len(reminders) == 1
        and "tps" in _text(reminders[0].get("text"))
        and reminders[0].get("date") == "2026-07-24"
        and reminders[0].get("time") == "15:00"
        and len(messages) == 1
        and "casey" in _text(messages[0].get("to"))
        and "tps" in _text(messages[0].get("text"))
        and any(term in _text(messages[0].get("text")) for term in ("friday", "end of day", "eod"))
    )
    return _base(valid, True, _effects(actions, ["set_reminder", "send_message"]), not artifacts, _state_exact(state, messages=1, reminders=1))


def learn_store(evidence):
    state, actions, artifacts = _inputs(evidence)
    saved = _successful(actions, "save_memory")
    valid_facts = [
        fact for fact in evidence.memory
        if "25" in _text(fact) and any(term in _text(fact) for term in ("10am", "10 am", "10:00"))
    ]
    saved_facts = [_text(item["args"].get("fact")) for item in saved]
    saved_correctly = (
        len(saved_facts) == 1
        and "25" in saved_facts[0]
        and any(term in saved_facts[0] for term in ("10am", "10 am", "10:00"))
    )
    valid = saved_correctly and len(valid_facts) == 1
    return _base(valid, True, _effects(actions, []), not artifacts, _state_exact(state))


def learn_use(evidence):
    state, actions, artifacts = _inputs(evidence)
    events = _new_events(state)
    valid = len(events) == 1
    if valid:
        event = events[0]
        start, end = _minutes(event["start"]), _minutes(event["end"])
        valid = (
            "priya" in _text(event.get("title")) + " " + _text(" ".join(event.get("attendees", [])))
            and event.get("date") == "2026-07-21"
            and 600 <= start <= 719
            and end - start == 25
        )
    return _base(valid, True, _effects(actions, ["add_event"]), not artifacts, _state_exact(state, events=8))


def multi_offsite(evidence):
    state, actions, artifacts = _inputs(evidence)
    events, sent = _new_events(state), state["sent_emails"]
    event_ok = len(events) == 1 and (
        _text(events[0].get("title")) == "summer offsite"
        and events[0].get("date") == "2026-07-24"
        and events[0].get("start") == "09:00"
        and events[0].get("end") == "16:00"
        and "lakeside pavilion" in _text(events[0].get("location"))
    )
    email_ok = False
    if len(sent) == 1:
        blob = _text(sent[0].get("subject")) + " " + _text(sent[0].get("body"))
        email_ok = _text(sent[0].get("to")) == "ceo@corp.com" and any(
            term in blob for term in ("confirm", "attend", "be there", "join", "count me in")
        )
    slide_ok = False
    if "offsite.pptx" in artifacts:
        slides = _slides(artifacts["offsite.pptx"])
        if len(slides) == 1:
            bullets = _text(" ".join(slides[0][1]))
            slide_ok = (
                _text(slides[0][0]) == "summer offsite"
                and ("july 24" in bullets or "2026-07-24" in bullets)
                and "09:00" in bullets
                and "16:00" in bullets
                and "lakeside pavilion" in bullets
            )
    return _base(
        event_ok and email_ok and slide_ok,
        _read(actions, "read_email", lambda args: args.get("id") == "e8"),
        _effects(actions, ["add_event", "send_email", "create_presentation"]),
        set(artifacts) == {"offsite.pptx"},
        _state_exact(state, events=8, emails=1),
    )


_EVALUATORS = {
    "pptx_basic": pptx_basic,
    "pptx_from_email": pptx_from_email,
    "xlsx_basic": xlsx_basic,
    "xlsx_from_email": xlsx_from_email,
    "email_reply": email_reply,
    "cal_add": cal_add,
    "cal_freeslot": cal_freeslot,
    "cal_brief": cal_brief,
    "remind_msg": remind_msg,
    "learn_store": learn_store,
    "learn_use": learn_use,
    "multi_offsite": multi_offsite,
}

GRADERS = {
    task_id: GraderSpec(
        id=f"office_demo.{task_id}",
        version=GRADER_VERSION,
        checks=CHECKS,
        evaluate=evaluator,
    )
    for task_id, evaluator in _EVALUATORS.items()
}
