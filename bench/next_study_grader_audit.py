"""Separate human-ground-truth and full-suite machine grader conformance."""

import argparse
import copy
from io import BytesIO
import json
import os
from pathlib import Path
import re
import stat
import tempfile

from openpyxl import load_workbook
from pptx import Presentation

from domains.office_demo import office_files
from domains.office_demo.reviewed_grader_v2 import (
    GRADER_IDENTITY, build_grader, task_id_for,
)
from harness.evidence import canonical_json_bytes
from harness.grading import GradingEvidence
from harness.instances import load_canonical_json, sha256_bytes, validate_manifest

from .next_study_review import (
    digest_review_artifact, review_packet, validate_adjudicated_outcomes,
)
from .next_study_validated_outcomes import validate_validated_outcomes


AUDIT_SCHEMA = "brick.next-study.grader-human-ground-truth-audit/1"
MACHINE_CONFORMANCE_SCHEMA = "brick.next-study.grader-machine-conformance/1"
VALIDATED_CONFORMANCE_SCHEMA = "brick.next-study.grader-validated-conformance/1"
_BUSINESS = frozenset((
    "create_presentation", "create_spreadsheet", "send_email", "add_event",
    "send_message", "set_reminder", "save_memory",
))
_SOURCES = frozenset(("source_read", "sources_read", "calendar_read"))


class NextStudyGraderAuditError(ValueError):
    pass


def _action(actions, tool, args):
    actions.append({
        "tool": tool, "args": copy.deepcopy(args), "ok": True,
        "result": "adjudicated positive baseline",
    })


def _column_name(index):
    value, result = index + 1, ""
    while value:
        value, remainder = divmod(value - 1, 26)
        result = chr(ord("A") + remainder) + result
    return result


def build_positive_evidence(packet, adjudicated_outcome, workdir):
    """Construct conformance evidence from the reviewed outcome, not an oracle."""

    state = copy.deepcopy(packet["initial_state"])
    memory = list(state.pop("memory"))
    state.pop("artifacts")
    actions = []
    artifacts_dir = Path(workdir) / "artifacts"
    artifacts_dir.mkdir(parents=True, exist_ok=True)
    for effect in adjudicated_outcome["outcome"]:
        kind = effect["type"]
        if kind == "source_read":
            if effect.get("list_required"):
                _action(actions, "list_emails", {})
            _action(actions, "read_email", {"id": effect["id"]})
        elif kind == "sources_read":
            if effect.get("list_required"):
                _action(actions, "list_emails", {})
            for identifier in effect["ids"]:
                _action(actions, "read_email", {"id": identifier})
        elif kind == "calendar_read":
            _action(actions, "list_events", {"date": effect["date"]})
        elif kind == "presentation_created":
            values = list(effect.get("required_values", []))
            values_by_slide = effect.get(
                "required_values_by_slide",
                [[] for _ in range(effect["exact_slide_count"])],
            )
            minimums = effect.get(
                "minimum_bullets_by_slide", [0] * effect["exact_slide_count"]
            )
            slides = []
            for index, title in enumerate(effect["ordered_titles"]):
                bullets = [str(value) for value in values_by_slide[index]]
                if values_by_slide is None and len(values) == effect["exact_slide_count"] - 1 and index:
                    bullets.append(str(values[index - 1]))
                elif values_by_slide is None and effect["exact_slide_count"] == 1:
                    bullets.extend(str(value) for value in values)
                while len(bullets) < minimums[index]:
                    bullets.append("Verified detail %d" % (len(bullets) + 1))
                slides.append({"title": title, "bullets": bullets})
            args = {"filename": effect["filename"], "slides": slides}
            office_files.create_presentation(str(artifacts_dir), **args)
            _action(actions, "create_presentation", args)
        elif kind == "spreadsheet_created":
            wanted, cents = effect.get("ordered_rows"), False
            if wanted is None:
                wanted, cents = effect["ordered_rows_cents"], True
            rows = [[str(value) for value in effect["headers"]]]
            for source_row in wanted:
                row = []
                for index, value in enumerate(source_row):
                    row.append(
                        "%d.%02d" % (value // 100, value % 100)
                        if cents and index == len(source_row) - 1 else str(value)
                    )
                rows.append(row)
            total = [""] * len(effect["headers"])
            total[0] = "Total"
            column = _column_name(len(total) - 1)
            total[-1] = "=SUM(%s2:%s%d)" % (column, column, len(rows))
            rows.append(total)
            args = {"filename": effect["filename"], "rows": rows}
            office_files.create_spreadsheet(str(artifacts_dir), **args)
            _action(actions, "create_spreadsheet", args)
        elif kind == "email_sent":
            record = {
                "to": effect["to"],
                "subject": "Re: %s" % effect.get("subject_contains", "attendance"),
                "body": "I confirm that I will attend. Count me in. %s" % " ".join(
                    str(value) for value in effect.get("required_mentions", [])
                ),
            }
            state["sent_emails"].append(record)
            _action(actions, "send_email", record)
        elif kind == "event_created":
            args = {
                "title": effect["title"], "date": effect["date"],
                "start_time": effect["start"], "end_time": effect["end"],
                "attendees": list(effect["attendees"]),
                "location": effect.get("location", ""),
            }
            state["events"].append({
                "id": "c%d" % (len(state["events"]) + 1),
                "title": args["title"], "date": args["date"],
                "start": args["start_time"], "end": args["end_time"],
                "attendees": args["attendees"], "location": args["location"],
            })
            _action(actions, "add_event", args)
        elif kind == "message_sent":
            mentions = effect.get("ordered_mentions", effect.get("required_mentions", []))
            parts = []
            for mention in mentions:
                matching = [item for item in state["events"] if item["title"] == mention]
                parts.append(
                    "%s at %s" % (mention, matching[0]["start"])
                    if effect.get("include_start_times") and matching else str(mention)
                )
            if effect.get("body_intent") == "deadline_commitment":
                parts.append(
                    "I will complete the full checklist by %s."
                    % effect.get("deadline", "the deadline")
                )
            record = {"to": effect["to"], "text": "; ".join(parts)}
            state["messages"].append(record)
            _action(actions, "send_message", record)
        elif kind == "reminder_created":
            record = {
                "text": "; ".join(effect["required_mentions"]),
                "date": effect["date"], "time": effect["time"],
            }
            state["reminders"].append(record)
            _action(actions, "set_reminder", record)
        elif kind == "memory_saved":
            fact = "; ".join(effect["required_facts"])
            memory.append(fact)
            _action(actions, "save_memory", {"fact": fact})
        else:
            raise NextStudyGraderAuditError("unsupported adjudicated effect %r" % kind)
    artifacts = [
        (path.name, path.read_bytes()) for path in sorted(artifacts_dir.iterdir())
        if path.is_file()
    ]
    return GradingEvidence.from_values(
        domain="office_demo", domain_version="0.1.0",
        task_id=task_id_for(packet, adjudicated_outcome), state=state,
        actions=actions, memory=memory, artifacts=artifacts,
    )


def _rebuild(evidence, *, state=None, actions=None, memory=None, artifacts=None):
    return GradingEvidence.from_values(
        domain=evidence.domain, domain_version=evidence.domain_version,
        task_id=evidence.task_id, state=evidence.state if state is None else state,
        actions=evidence.actions if actions is None else actions,
        memory=evidence.memory if memory is None else memory,
        artifacts=(sorted(evidence.artifact_map().items()) if artifacts is None else artifacts),
    )


def _wrong_artifact(name, payload):
    if name.endswith(".pptx"):
        deck = Presentation(BytesIO(payload))
        deck.slides[0].shapes.title.text = "Intentionally incorrect title"
        output = BytesIO()
        deck.save(output)
        return output.getvalue()
    if name.endswith(".xlsx"):
        book = load_workbook(BytesIO(payload))
        book.worksheets[0].cell(row=2, column=book.worksheets[0].max_column).value = 0
        output = BytesIO()
        book.save(output)
        book.close()
        return output.getvalue()
    raise NextStudyGraderAuditError("unexpected artifact type")


def _required_mutation(evidence, effects):
    artifacts = sorted(evidence.artifact_map().items())
    if artifacts:
        name, payload = artifacts[0]
        return _rebuild(evidence, artifacts=[(name, _wrong_artifact(name, payload))] + artifacts[1:])
    state, memory = copy.deepcopy(evidence.state), list(evidence.memory)
    kinds = {item["type"] for item in effects}
    if "email_sent" in kinds:
        state["sent_emails"][-1]["body"] = "Intentionally incorrect"
    elif "event_created" in kinds:
        state["events"][-1]["title"] = "Intentionally incorrect"
    elif "message_sent" in kinds:
        state["messages"][-1]["text"] = "Intentionally incorrect"
    elif "reminder_created" in kinds:
        state["reminders"][-1]["text"] = "Intentionally incorrect"
    elif "memory_saved" in kinds:
        memory[-1] = "Intentionally incorrect"
    else:
        raise NextStudyGraderAuditError("case has no mutable required outcome")
    return _rebuild(evidence, state=state, memory=memory)


def _mutations(evidence, effects):
    actions = list(evidence.actions)
    business_index = next(
        index for index, item in enumerate(actions)
        if item["ok"] and item["tool"] in _BUSINESS
    )
    extra_state = copy.deepcopy(evidence.state)
    extra_state["messages"].append({"to": "nobody", "text": "unexpected"})
    result = {
        "required_outcome": _required_mutation(evidence, effects),
        "exact_business_effects": _rebuild(
            evidence, actions=actions[:business_index] + actions[business_index + 1:]
        ),
        "exact_artifacts": _rebuild(
            evidence, artifacts=sorted(evidence.artifact_map().items())
            + [("unexpected.bin", b"unexpected")]
        ),
        "no_unrequested_state": _rebuild(evidence, state=extra_state),
        "no_unauthorized_effects": _rebuild(evidence, actions=actions + [{
            "tool": "unauthorized_probe", "args": {}, "ok": True,
            "result": "must fail closed",
        }]),
    }
    if any(item["type"] in _SOURCES for item in effects):
        result["source_observed"] = _rebuild(
            evidence,
            actions=[item for item in actions if item["tool"] not in {"read_email", "list_events"}],
        )
    return result


def _replace_artifact(evidence, filename, payload):
    return _rebuild(evidence, artifacts=[
        (name, payload if name == filename else value)
        for name, value in sorted(evidence.artifact_map().items())
    ])


def _mutate_presentation_value(payload, slide_index, wanted, occurrence):
    deck = Presentation(BytesIO(payload))
    target = str(wanted).strip().casefold()
    matches = [
        paragraph
        for shape in deck.slides[slide_index].shapes
        if getattr(shape, "has_text_frame", False)
        for paragraph in shape.text_frame.paragraphs
        if paragraph.text.strip().casefold() == target
    ]
    if occurrence < len(matches):
        matches[occurrence].text = matches[occurrence].text + "__wrong"
        output = BytesIO(); deck.save(output)
        return output.getvalue()
    raise NextStudyGraderAuditError("graded presentation fact is absent from baseline")


def _reverse_presentation_values(payload, slide_index, required):
    deck = Presentation(BytesIO(payload))
    wanted = [str(value).strip().casefold() for value in required]
    paragraphs = [
        paragraph
        for shape in deck.slides[slide_index].shapes
        if getattr(shape, "has_text_frame", False)
        for paragraph in shape.text_frame.paragraphs
    ]
    matches = [paragraph for paragraph in paragraphs if paragraph.text.strip().casefold() in wanted]
    if [paragraph.text.strip().casefold() for paragraph in matches] == wanted:
        values = [paragraph.text for paragraph in matches]
        for paragraph, value in zip(matches, reversed(values)):
            paragraph.text = value
        output = BytesIO(); deck.save(output)
        return output.getvalue()
    raise NextStudyGraderAuditError("ordered presentation facts are absent from baseline")


def _semantic_mutations(evidence, effects):
    """Known-bad semantic outputs that every strict grader must reject."""

    probes = []
    for effect_index, effect in enumerate(effects):
        kind = effect["type"]
        if kind == "sources_read" and effect.get("list_required"):
            probes.append((
                "source_list_%d" % effect_index, "source_observed",
                _rebuild(evidence, actions=[
                    item for item in evidence.actions if item["tool"] != "list_emails"
                ]),
            ))
        elif kind == "email_sent":
            state = copy.deepcopy(evidence.state)
            match = next(item for item in state["sent_emails"] if item["to"] == effect["to"])
            match["body"] = "I cannot confirm and will not attend. " + " ".join(
                effect.get("required_mentions", [])
            )
            probes.append(("negated_email_%d" % effect_index, "required_outcome", _rebuild(evidence, state=state)))
            dates = [
                value for value in effect.get("required_mentions", [])
                if re.fullmatch(r"\d{4}-\d{2}-\d{2}", str(value))
            ]
            if dates:
                state = copy.deepcopy(evidence.state)
                match = next(item for item in state["sent_emails"] if item["to"] == effect["to"])
                match["body"] = match["body"].replace(dates[0], "February 5, 2028")
                probes.append((
                    "reformatted_email_field_%d" % effect_index,
                    "required_outcome", _rebuild(evidence, state=state),
                ))
        elif kind == "message_sent" and effect.get("body_intent") == "deadline_commitment":
            state = copy.deepcopy(evidence.state)
            match = next(item for item in state["messages"] if item["to"] == effect["to"])
            match["text"] = "; ".join(effect.get("required_mentions", [])) + "; will not be complete by the deadline"
            probes.append(("negated_deadline_%d" % effect_index, "required_outcome", _rebuild(evidence, state=state)))
            if effect.get("deadline"):
                state = copy.deepcopy(evidence.state)
                match = next(item for item in state["messages"] if item["to"] == effect["to"])
                match["text"] = match["text"].replace(effect["deadline"], "the deadline")
                probes.append((
                    "missing_exact_deadline_%d" % effect_index,
                    "required_outcome", _rebuild(evidence, state=state),
                ))
        elif kind == "message_sent" and effect.get("forbid_date_tokens"):
            state = copy.deepcopy(evidence.state)
            match = next(item for item in state["messages"] if item["to"] == effect["to"])
            match["text"] += "; 2099-12-31"
            probes.append(("forbidden_date_%d" % effect_index, "required_outcome", _rebuild(evidence, state=state)))
        if kind == "message_sent" and effect.get("forbidden_mentions"):
            state = copy.deepcopy(evidence.state)
            match = next(item for item in state["messages"] if item["to"] == effect["to"])
            match["text"] += "; " + effect["forbidden_mentions"][0]
            probes.append((
                "forbidden_mention_%d" % effect_index,
                "required_outcome", _rebuild(evidence, state=state),
            ))
        if kind in {"message_sent", "reminder_created"} and effect.get("exact_mentions"):
            state = copy.deepcopy(evidence.state)
            prefix = re.sub(r"[0-9]+$", "", effect["required_mentions"][0])
            extra = prefix + "999"
            if kind == "message_sent":
                match = next(item for item in state["messages"] if item["to"] == effect["to"])
            else:
                match = next(
                    item for item in state["reminders"]
                    if item["date"] == effect["date"] and item["time"] == effect["time"]
                )
            match["text"] += "; " + extra
            probes.append((
                "extra_identifier_%d" % effect_index,
                "required_outcome", _rebuild(evidence, state=state),
            ))
        elif kind == "memory_saved":
            memory = list(evidence.memory)
            memory[-1] += "; unrequested_conflict=true"
            probes.append(("memory_exactness_%d" % effect_index, "required_outcome", _rebuild(evidence, memory=memory)))
        elif kind == "presentation_created":
            payload = evidence.artifact_map()[effect["filename"]]
            for slide_index, values in enumerate(effect.get("required_values_by_slide", [])):
                seen = {}
                for value_index, value in enumerate(values):
                    key = str(value).strip().casefold()
                    occurrence = seen.get(key, 0)
                    seen[key] = occurrence + 1
                    probes.append((
                        "presentation_fact_%d_%d_%d" % (effect_index, slide_index, value_index),
                        "required_outcome",
                        _replace_artifact(
                            evidence, effect["filename"],
                            _mutate_presentation_value(
                                payload, slide_index, value, occurrence,
                            ),
                        ),
                    ))
                if len(values) > 1:
                    probes.append((
                        "presentation_order_%d_%d" % (effect_index, slide_index),
                        "required_outcome",
                        _replace_artifact(
                            evidence, effect["filename"],
                            _reverse_presentation_values(payload, slide_index, values),
                        ),
                    ))
    return probes


_SEMANTIC_PROBE_PREFIXES = (
    "source_list", "negated_email", "negated_deadline", "forbidden_date",
    "reformatted_email_field", "missing_exact_deadline", "forbidden_mention",
    "extra_identifier", "memory_exactness", "presentation_fact",
    "presentation_order",
)


def _semantic_probe_inventory(outcomes):
    """Derive the required semantic probes without consulting executed probes."""

    counts = {key: 0 for key in _SEMANTIC_PROBE_PREFIXES}
    for record in outcomes.values():
        for effect in record["outcome"]:
            kind = effect["type"]
            if kind in {"source_read", "sources_read"} and effect.get("list_required"):
                counts["source_list"] += 1
            if kind == "email_sent":
                counts["negated_email"] += 1
                if any(
                    re.fullmatch(r"\d{4}-\d{2}-\d{2}", str(value))
                    for value in effect.get("required_mentions", [])
                ):
                    counts["reformatted_email_field"] += 1
            if kind == "message_sent" and effect.get("body_intent") == "deadline_commitment":
                counts["negated_deadline"] += 1
                if effect.get("deadline"):
                    counts["missing_exact_deadline"] += 1
            if kind == "message_sent" and effect.get("forbid_date_tokens"):
                counts["forbidden_date"] += 1
            if kind == "message_sent" and effect.get("forbidden_mentions"):
                counts["forbidden_mention"] += 1
            if kind in {"message_sent", "reminder_created"} and effect.get("exact_mentions"):
                counts["extra_identifier"] += 1
            if kind == "memory_saved":
                counts["memory_exactness"] += 1
            if kind == "presentation_created":
                values_by_slide = effect.get("required_values_by_slide", [])
                counts["presentation_fact"] += sum(len(values) for values in values_by_slide)
                counts["presentation_order"] += sum(
                    len(values) > 1 for values in values_by_slide
                )
    return counts


def _manifest_instances(manifests):
    instances = []
    for manifest in manifests:
        validate_manifest(manifest)
        instances.extend(manifest["instances"])
    if len(instances) != 528:
        raise NextStudyGraderAuditError("grader conformance requires the frozen 528 cases")
    return instances


def _audit_records(manifests, outcome_records, *, schema_version, evidence_basis,
                   source_artifact_sha256, may_satisfy_human_gate):
    instances = _manifest_instances(manifests)
    by_id = {item["content"]["id"]: item for item in instances}
    outcomes = {item["instance_id"]: item for item in outcome_records}
    if (
        len(outcomes) != len(outcome_records)
        or not outcomes or not set(outcomes) <= set(by_id)
    ):
        raise NextStudyGraderAuditError("grader audit outcomes are duplicate or unknown")
    probe_count, check_counts = 0, {}
    semantic_probe_counts = {key: 0 for key in _SEMANTIC_PROBE_PREFIXES}
    benign_counts = {
        "equivalent_serialization": 0,
        "failed_unauthorized_call": 0,
        "repeated_safe_read": 0,
        "nonbusiness_scratch_memory": 0,
    }
    with tempfile.TemporaryDirectory(prefix="brick-next-grader-audit-") as root:
        for index, instance_id in enumerate(sorted(outcomes)):
            packet, outcome = review_packet(by_id[instance_id]), outcomes[instance_id]
            grader = build_grader(packet, outcome)
            evidence = build_positive_evidence(packet, outcome, Path(root) / str(index))
            if grader.grade_evidence(evidence).strict_success is not True:
                raise NextStudyGraderAuditError("positive baseline failed %s" % instance_id)
            failed_call = _rebuild(evidence, actions=list(evidence.actions) + [{
                "tool": "unauthorized_probe", "args": {}, "ok": False,
                "result": "benign rejected call",
            }])
            if grader.grade_evidence(failed_call).strict_success is not True:
                raise NextStudyGraderAuditError("benign control rejected %s" % instance_id)
            benign_counts["failed_unauthorized_call"] += 1
            equivalent = _rebuild(
                evidence,
                state=dict(reversed(list(evidence.state.items()))),
                actions=[dict(reversed(list(item.items()))) for item in evidence.actions],
                artifacts=list(reversed(sorted(evidence.artifact_map().items()))),
            )
            if grader.grade_evidence(equivalent).strict_success is not True:
                raise NextStudyGraderAuditError(
                    "equivalent serialization rejected %s" % instance_id
                )
            benign_counts["equivalent_serialization"] += 1
            safe_reads = [
                item for item in evidence.actions
                if item["ok"] and item["tool"] in {"read_email", "list_events"}
            ]
            if safe_reads:
                repeated = _rebuild(
                    evidence, actions=list(evidence.actions) + [copy.deepcopy(safe_reads[0])]
                )
                if grader.grade_evidence(repeated).strict_success is not True:
                    raise NextStudyGraderAuditError(
                        "repeated safe read rejected %s" % instance_id
                    )
                benign_counts["repeated_safe_read"] += 1
            if not any(item["type"] == "memory_saved" for item in outcome["outcome"]):
                scratch = "nonbusiness conformance scratch"
                with_scratch = _rebuild(
                    evidence,
                    memory=list(evidence.memory) + [scratch],
                    actions=list(evidence.actions) + [{
                        "tool": "save_memory", "args": {"fact": scratch},
                        "ok": True, "result": "benign internal scratch",
                    }],
                )
                if grader.grade_evidence(with_scratch).strict_success is not True:
                    raise NextStudyGraderAuditError(
                        "nonbusiness scratch memory biased grading %s" % instance_id
                    )
                benign_counts["nonbusiness_scratch_memory"] += 1
            generic_mutations = [
                (check_id, check_id, mutated)
                for check_id, mutated in sorted(_mutations(evidence, outcome["outcome"]).items())
            ]
            semantic_mutations = _semantic_mutations(evidence, outcome["outcome"])
            semantic_ids = {item[0] for item in semantic_mutations}
            mutations = generic_mutations + semantic_mutations
            for _probe_id, check_id, mutated in mutations:
                graded = grader.grade_evidence(mutated)
                checks = {key: value for key, _description, value in graded.checks}
                if graded.grader_status != "graded" or graded.strict_success is not False:
                    raise NextStudyGraderAuditError("mutation survived %s" % instance_id)
                if checks.get(check_id) is not False:
                    raise NextStudyGraderAuditError("target check survived %s" % instance_id)
                probe_count += 1
                check_counts[check_id] = check_counts.get(check_id, 0) + 1
                if _probe_id in semantic_ids:
                    prefix = next(
                        key for key in _SEMANTIC_PROBE_PREFIXES
                        if _probe_id.startswith(key + "_")
                    )
                    semantic_probe_counts[prefix] += 1
    source_cases = sum(
        any(item["type"] in _SOURCES for item in outcome["outcome"])
        for outcome in outcomes.values()
    )
    case_count = len(outcomes)
    expected_semantic = _semantic_probe_inventory(outcomes)
    expected_generic = 5 * case_count + source_cases
    expected_probes = expected_generic + sum(expected_semantic.values())
    expected_benign = {
        "equivalent_serialization": case_count,
        "failed_unauthorized_call": case_count,
        "repeated_safe_read": source_cases,
        "nonbusiness_scratch_memory": sum(
            not any(item["type"] == "memory_saved" for item in outcome["outcome"])
            for outcome in outcomes.values()
        ),
    }
    if probe_count != expected_probes or semantic_probe_counts != expected_semantic:
        raise NextStudyGraderAuditError("targeted mutation inventory drifted")
    if benign_counts != expected_benign:
        raise NextStudyGraderAuditError("mutation or benign control count drifted")
    grader_path = Path(__file__).resolve().parents[1] / "domains" / "office_demo" / "reviewed_grader_v2.py"
    return {
        "schema_version": schema_version,
        "evidence_basis": evidence_basis,
        "may_satisfy_human_ground_truth_gate": may_satisfy_human_gate,
        "grader_identity": GRADER_IDENTITY,
        "grader_source_sha256": sha256_bytes(grader_path.read_bytes()),
        "source_artifact_sha256": source_artifact_sha256,
        "manifest_set_sha256": digest_review_artifact(
            sorted(manifests, key=lambda item: item["split"])
        ),
        "case_count": case_count,
        "positive_baselines": case_count,
        "targeted_mutations": expected_probes,
        "generic_mutation_count": expected_generic,
        "semantic_probe_counts": dict(sorted(semantic_probe_counts.items())),
        "benign_non_rejection_controls": sum(benign_counts.values()),
        "benign_control_counts": benign_counts,
        "check_probe_counts": dict(sorted(check_counts.items())),
        "passed": True,
        "live_model_calls": 0,
    }


def audit_all(manifests, adjudicated_outcomes):
    """Audit the research grader against 308 independently reviewed outcomes."""

    try:
        validate_adjudicated_outcomes(adjudicated_outcomes, manifests)
    except ValueError as exc:
        raise NextStudyGraderAuditError(str(exc))
    return _audit_records(
        manifests, adjudicated_outcomes["records"], schema_version=AUDIT_SCHEMA,
        evidence_basis="sealed_human_validity_outcomes",
        source_artifact_sha256=digest_review_artifact(adjudicated_outcomes),
        may_satisfy_human_gate=True,
    )


def audit_machine_conformance(manifests):
    """Exercise all 528 cases using the independent parser-oracle as fixtures."""

    from domains.office_demo.outcome_oracle_v2 import derive_outcome

    records = []
    for instance in _manifest_instances(manifests):
        packet = review_packet(instance)
        records.append({
            "instance_id": instance["content"]["id"],
            "content_sha256": instance["content_sha256"],
            "review_packet_sha256": digest_review_artifact(packet),
            "prompt_valid": True,
            "outcome": derive_outcome(
                packet["family"], packet["prompt"], packet["subepisode_prompts"],
                packet["initial_state"], packet["today"],
            ),
            "accepted_alternatives": [],
            "review_resolution": "machine_oracle_fixture",
        })
    oracle_path = Path(__file__).resolve().parents[1] / "domains" / "office_demo" / "outcome_oracle_v2.py"
    return _audit_records(
        manifests, records, schema_version=MACHINE_CONFORMANCE_SCHEMA,
        evidence_basis="independent_parser_oracle_fixture_not_human_review",
        source_artifact_sha256=sha256_bytes(oracle_path.read_bytes()),
        may_satisfy_human_gate=False,
    )


def audit_validated_conformance(manifests, validated_outcomes):
    """Exercise all cases from the independently compiled public outcomes."""

    validate_validated_outcomes(validated_outcomes, manifests)
    return _audit_records(
        manifests,
        validated_outcomes["records"],
        schema_version=VALIDATED_CONFORMANCE_SCHEMA,
        evidence_basis="independently_compiled_public_packet_outcomes",
        source_artifact_sha256=digest_review_artifact(validated_outcomes),
        may_satisfy_human_gate=False,
    )


def _publish_marker_last(path, document):
    path = Path(path)
    marker = Path(str(path) + ".complete")
    path.parent.mkdir(parents=True, exist_ok=True)
    if path.exists() or marker.exists():
        raise NextStudyGraderAuditError("grader audit publication already exists")
    descriptor = os.open(str(path), os.O_WRONLY | os.O_CREAT | os.O_EXCL, 0o600)
    try:
        with os.fdopen(descriptor, "wb") as stream:
            stream.write(canonical_json_bytes(document, allow_float=False, newline=True))
            stream.flush()
            os.fsync(stream.fileno())
    except Exception:
        raise
    marker_descriptor = os.open(
        str(marker), os.O_WRONLY | os.O_CREAT | os.O_EXCL, 0o600
    )
    try:
        os.fsync(marker_descriptor)
    finally:
        os.close(marker_descriptor)
    return path


def _load_marker_last(path):
    path = Path(path)
    marker = Path(str(path) + ".complete")
    try:
        file_stat = os.lstat(path)
        marker_stat = os.lstat(marker)
    except FileNotFoundError:
        raise NextStudyGraderAuditError(
            "adjudicated outcomes or completion marker is missing"
        )
    if (
        not stat.S_ISREG(file_stat.st_mode)
        or not stat.S_ISREG(marker_stat.st_mode)
        or marker_stat.st_size != 0
    ):
        raise NextStudyGraderAuditError("adjudicated outcomes are not marker-last")
    return load_canonical_json(path)


def main(argv=None):
    root = Path(__file__).resolve().parents[1]
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--outcomes", type=Path,
        default=None,
    )
    parser.add_argument(
        "--output", type=Path,
    )
    modes = parser.add_mutually_exclusive_group()
    modes.add_argument(
        "--machine-conformance", action="store_true",
        help="Run all 528 parser-oracle fixtures as advisory conformance",
    )
    modes.add_argument(
        "--validated-conformance", action="store_true",
        help="Run all 528 independently compiled public-packet outcomes",
    )
    args = parser.parse_args(argv)
    manifests = [
        load_canonical_json(root / "bench" / "manifests" / "office-v2" / (split + ".json"))
        for split in (
            "development", "calibration", "validation", "sentinel",
            "retained", "adversarial",
        )
    ]
    if args.machine_conformance:
        result = audit_machine_conformance(manifests)
    elif args.validated_conformance:
        outcome_path = args.outcomes or (
            root / "evidence" / "next-study"
            / "office-v2-validated-outcomes.json"
        )
        result = audit_validated_conformance(
            manifests, load_canonical_json(outcome_path)
        )
    else:
        outcome_path = args.outcomes or (
            root / "evidence" / "next-study"
            / "office-v2-adjudicated-outcomes.json"
        )
        result = audit_all(manifests, _load_marker_last(outcome_path))
    output = args.output or root / "evidence" / "next-study" / (
        "office-v2-grader-machine-conformance.json"
        if args.machine_conformance or args.validated_conformance
        else "office-v2-grader-human-ground-truth-audit.json"
    )
    _publish_marker_last(output, result)
    print(json.dumps({
        "status": "passed", "cases": result["case_count"],
        "targeted_mutations": result["targeted_mutations"],
        "benign_controls": result["benign_non_rejection_controls"],
        "live_model_calls": 0, "output": str(output),
    }, sort_keys=True))


if __name__ == "__main__":
    main()


__all__ = [
    "AUDIT_SCHEMA", "MACHINE_CONFORMANCE_SCHEMA", "NextStudyGraderAuditError",
    "audit_all", "audit_machine_conformance", "audit_validated_conformance",
    "build_positive_evidence", "main",
]
