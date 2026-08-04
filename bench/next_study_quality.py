"""Offline outcome-validity probes for generated office graders.

The probes are model-free.  For every frozen generated case, they construct a
known-good rules-reference result and then apply one targeted, schema-valid
mutation per applicable named grader check.  A probe passes only when the named
check and strict success both become false without a grader error.
"""

import copy
from io import BytesIO
from pathlib import Path
import tempfile

from openpyxl import load_workbook
from pptx import Presentation

from domains.office_demo.generated_grader import build_grader
from domains.office_demo.rules_reference import execute as execute_rules
from harness.grading import GradingEvidence
from harness.instances import load_canonical_json, validate_manifest


ROOT = Path(__file__).resolve().parents[1]
MANIFESTS = ROOT / "bench" / "manifests" / "office-v1"

_BUSINESS = frozenset(
    {
        "create_presentation",
        "create_spreadsheet",
        "send_email",
        "add_event",
        "send_message",
        "set_reminder",
        "save_memory",
    }
)
_SOURCE_EFFECTS = frozenset({"source_read", "sources_read", "calendar_read"})


def _rebuild(evidence, *, state=None, actions=None, memory=None, artifacts=None):
    return GradingEvidence.from_values(
        domain=evidence.domain,
        domain_version=evidence.domain_version,
        task_id=evidence.task_id,
        state=evidence.state if state is None else state,
        actions=evidence.actions if actions is None else actions,
        memory=evidence.memory if memory is None else memory,
        artifacts=(
            sorted(evidence.artifact_map().items())
            if artifacts is None
            else artifacts
        ),
    )


def _wrong_artifact(name, payload):
    if name.endswith(".pptx"):
        deck = Presentation(BytesIO(payload))
        deck.slides[0].shapes.title.text = "Intentionally incorrect title"
        output = BytesIO()
        deck.save(output)
        return output.getvalue()
    if name.endswith(".xlsx"):
        book = load_workbook(BytesIO(payload))
        sheet = book.worksheets[0]
        sheet.cell(row=2, column=sheet.max_column).value = 0
        output = BytesIO()
        book.save(output)
        book.close()
        return output.getvalue()
    raise AssertionError("unexpected generated artifact %r" % name)


def _required_outcome_mutation(evidence, content):
    artifacts = sorted(evidence.artifact_map().items())
    if artifacts:
        name, payload = artifacts[0]
        return _rebuild(
            evidence,
            artifacts=[(name, _wrong_artifact(name, payload))] + artifacts[1:],
        )
    state = copy.deepcopy(evidence.state)
    memory = list(evidence.memory)
    kinds = {effect["type"] for effect in content["required_effects"]}
    if "email_sent" in kinds:
        state["sent_emails"][-1]["body"] = "Intentionally incorrect"
    elif "event_created" in kinds:
        state["events"][-1]["title"] = "Intentionally incorrect"
    elif "message_sent" in kinds:
        state["messages"][-1]["text"] = "Intentionally incorrect"
    elif "reminder_created" in kinds:
        state["reminders"][-1]["text"] = "Intentionally incorrect"
    elif "memory_saved" in kinds:
        memory[-1] = "Intentionally incorrect"
    else:
        raise AssertionError("case has no mutable required outcome")
    return _rebuild(evidence, state=state, memory=memory)


def _mutations(evidence, content):
    actions = list(evidence.actions)
    first_business = next(
        index
        for index, item in enumerate(actions)
        if item["ok"] and item["tool"] in _BUSINESS
    )
    without_business = actions[:first_business] + actions[first_business + 1 :]
    extra_state = copy.deepcopy(evidence.state)
    extra_state["messages"].append({"to": "nobody", "text": "unexpected"})
    mutations = {
        "required_outcome": _required_outcome_mutation(evidence, content),
        "exact_business_effects": _rebuild(
            evidence,
            actions=without_business,
        ),
        "exact_artifacts": _rebuild(
            evidence,
            artifacts=sorted(evidence.artifact_map().items())
            + [("unexpected.bin", b"unexpected")],
        ),
        "no_unrequested_state": _rebuild(
            evidence,
            state=extra_state,
        ),
        "no_unauthorized_effects": _rebuild(
            evidence,
            actions=actions
            + [
                {
                    "tool": "unauthorized_probe",
                    "args": {},
                    "ok": True,
                    "result": "must fail closed",
                }
            ],
        ),
    }
    if any(
        effect["type"] in _SOURCE_EFFECTS
        for effect in content["required_effects"]
    ):
        mutations["source_observed"] = _rebuild(
            evidence,
            actions=[
                item
                for item in actions
                if item["tool"] not in {"read_email", "list_events"}
            ],
        )
    return mutations


def audit_instance(instance, workdir):
    content = instance["content"]
    evidence = execute_rules(instance, workdir)
    grader = build_grader(instance)
    baseline = grader.grade_evidence(evidence)
    if baseline.strict_success is not True:
        raise AssertionError("rules reference failed %s" % content["id"])
    passed = []
    for check_id, mutated in sorted(_mutations(evidence, content).items()):
        outcome = grader.grade_evidence(mutated)
        checks = {
            key: value for key, _description, value in outcome.checks
        }
        if outcome.grader_status != "graded":
            raise AssertionError(
                "%s mutation caused grader error for %s"
                % (check_id, content["id"])
            )
        if outcome.strict_success is not False or checks.get(check_id) is not False:
            raise AssertionError(
                "%s mutation survived for %s" % (check_id, content["id"])
            )
        passed.append(check_id)
    return passed


def audit_all(manifests_root=MANIFESTS):
    case_count = 0
    probe_count = 0
    check_counts = {}
    with tempfile.TemporaryDirectory(prefix="brick-grader-mutations-") as root:
        root_path = Path(root)
        for split in (
            "development",
            "validation",
            "sentinel",
            "retained",
            "adversarial",
        ):
            manifest = load_canonical_json(Path(manifests_root) / (split + ".json"))
            validate_manifest(manifest)
            for instance in manifest["instances"]:
                checks = audit_instance(
                    instance,
                    root_path / instance["content"]["id"].replace(".", "-"),
                )
                case_count += 1
                probe_count += len(checks)
                for check_id in checks:
                    check_counts[check_id] = check_counts.get(check_id, 0) + 1
    return {
        "schema_version": "brick.next-study.grader-mutation-audit/1",
        "generator_version": "office-generators/1.1.0",
        "case_count": case_count,
        "probe_count": probe_count,
        "check_probe_counts": dict(sorted(check_counts.items())),
        "all_applicable_mutations_rejected": case_count == 352,
        "live_model_calls": 0,
        "retained_model_execution": False,
    }


__all__ = ["audit_all", "audit_instance"]
