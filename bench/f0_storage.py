"""Disposable marker-last storage probe for the Lenovo F0 gate.

This module validates the filesystem primitives required by the planned S4
evidence store.  It is intentionally not the production benchmark store.
No directory is renamed or replaced: a candidate becomes visible only after
an empty ``COMMITTED`` marker is created exclusively and every prepared hash
validates.
"""

import argparse
import errno
import hashlib
import json
import math
import os
from pathlib import Path
import re
import subprocess
import sys
import tempfile
import threading
import time
import uuid

from openpyxl import Workbook
from pptx import Presentation


PREPARED_SCHEMA = "brick.attempt-prepared/1"
PREPARED = "PREPARED.json"
COMMITTED = "COMMITTED"
REQUIRED_FILES = (
    "actions.json",
    "artifacts/probe.pptx",
    "artifacts/probe.xlsx",
    "attempt.json",
    "final-state.json",
    "grade.json",
    "initial-state.json",
    "memory-delta.jsonl",
    "result.json",
    "transcript.md",
)
BOUNDARIES = (
    "candidate_created",
    "attempt_written",
    "initial_state_written",
    "final_state_written",
    "result_written",
    "grade_written",
    "actions_written",
    "transcript_written",
    "memory_written",
    "xlsx_written",
    "pptx_written",
    "prepared_written",
    "prepared_validated",
    "committed_created",
)
_COMPONENT = re.compile(r"^[a-zA-Z0-9][a-zA-Z0-9._-]{0,127}$")
_RETRYABLE_WINERRORS = frozenset({5, 32, 33})
_RETRYABLE_ERRNOS = frozenset({errno.EACCES, errno.EBUSY})
_DELAYS = (0.0, 0.05, 0.1, 0.2, 0.4, 0.8, 1.6)


class StorageIntegrityError(RuntimeError):
    """The candidate cannot be trusted or safely recovered."""


class InjectedStop(RuntimeError):
    """Test-only stand-in for a process exit at a write boundary."""


def _component(value, label):
    if not isinstance(value, str) or not _COMPONENT.fullmatch(value):
        raise ValueError(f"{label} is not a safe path component")
    return value


def _logical_hash(value):
    if not isinstance(value, str) or not re.fullmatch(
        r"[0-9a-f]{64}", value
    ):
        raise ValueError("logical id must be a lowercase SHA-256 digest")
    return value


def _physical_uuid(value):
    _component(value, "physical id")
    try:
        parsed = uuid.UUID(value)
    except (AttributeError, TypeError, ValueError) as exc:
        raise ValueError("physical id must be a UUID") from exc
    if str(parsed) != value:
        raise ValueError("physical id must use canonical UUID form")
    return value


def _sha256(path):
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _sync_parent(path):
    """Best-effort directory flush where the host exposes one."""
    flags = getattr(os, "O_DIRECTORY", 0) | os.O_RDONLY
    try:
        descriptor = os.open(str(path), flags)
    except OSError:
        return
    try:
        os.fsync(descriptor)
    except OSError:
        # Windows does not expose portable directory fsync semantics. Readers
        # still fail closed by validating marker and manifest after restart.
        pass
    finally:
        os.close(descriptor)


def _write_bytes(path, payload, exclusive=True):
    mode = "xb" if exclusive else "wb"
    with path.open(mode) as handle:
        handle.write(payload)
        handle.flush()
        os.fsync(handle.fileno())
    _sync_parent(path.parent)


def _write_json(path, value, exclusive=True):
    payload = (
        json.dumps(
            value,
            ensure_ascii=False,
            indent=2,
            sort_keys=True,
        )
        + "\n"
    ).encode("utf-8")
    _write_bytes(path, payload, exclusive=exclusive)


def _boundary(name, crash_after=None, hard_exit=False):
    if name != crash_after:
        return
    if hard_exit:
        os._exit(99)
    raise InjectedStop(name)


def create_office_templates(directory):
    """Create real, reopenable PPTX and XLSX payloads once per spike."""
    directory = Path(directory)
    directory.mkdir(parents=True, exist_ok=True)
    xlsx = directory / "template.xlsx"
    pptx = directory / "template.pptx"

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "F0"
    sheet.append(("probe", "value"))
    sheet.append(("marker-last", 1))
    workbook.save(str(xlsx))

    presentation = Presentation()
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[1]
    )
    slide.shapes.title.text = "Brick F0"
    slide.placeholders[1].text = "Marker-last Windows storage probe"
    presentation.save(str(pptx))

    # Reopen now so corrupt template generation fails before the spike.
    from openpyxl import load_workbook

    load_workbook(str(xlsx), read_only=True).close()
    Presentation(str(pptx))
    return {
        "artifacts/probe.xlsx": xlsx.read_bytes(),
        "artifacts/probe.pptx": pptx.read_bytes(),
    }


def candidate_path(root, logical_id, physical_id):
    return (
        Path(root)
        / "attempts"
        / _logical_hash(logical_id)
        / _physical_uuid(physical_id)
    )


def prepare_candidate(
    root,
    logical_id,
    physical_id,
    office_payloads,
    crash_after=None,
    hard_exit=False,
    commit=True,
    deadline_seconds=30.0,
    clock=time.monotonic,
    sleeper=time.sleep,
):
    """Create one unique candidate and optionally publish it."""
    candidate = candidate_path(root, logical_id, physical_id)
    candidate.parent.mkdir(parents=True, exist_ok=True)
    candidate.mkdir(exist_ok=False)
    _sync_parent(candidate.parent)
    _boundary("candidate_created", crash_after, hard_exit)

    attempt = {
        "schema_version": "brick.f0.attempt/1",
        "logical_hash": logical_id,
        "physical_id": physical_id,
        "probe": "lenovo_f0_marker_last",
    }
    _write_json(candidate / "attempt.json", attempt)
    _boundary("attempt_written", crash_after, hard_exit)

    _write_json(
        candidate / "initial-state.json",
        {"counter": 0, "schema_version": "brick.f0.state/1"},
    )
    _boundary("initial_state_written", crash_after, hard_exit)
    _write_json(
        candidate / "final-state.json",
        {"counter": 1, "schema_version": "brick.f0.state/1"},
    )
    _boundary("final_state_written", crash_after, hard_exit)
    _write_json(
        candidate / "result.json",
        {"execution_status": "done", "schema_version": "brick.f0.result/1"},
    )
    _boundary("result_written", crash_after, hard_exit)
    _write_json(
        candidate / "grade.json",
        {"schema_version": "brick.f0.grade/1", "strict_success": True},
    )
    _boundary("grade_written", crash_after, hard_exit)
    _write_json(
        candidate / "actions.json",
        {
            "actions": [{"tool": "f0_probe", "ok": True}],
            "schema_version": "brick.f0.actions/1",
        },
    )
    _boundary("actions_written", crash_after, hard_exit)
    _write_bytes(
        candidate / "transcript.md",
        b"# F0 storage probe\n",
    )
    _boundary("transcript_written", crash_after, hard_exit)
    _write_bytes(
        candidate / "memory-delta.jsonl",
        b'{"schema_version":"brick.f0.memory/1","delta":[]}\n',
    )
    _boundary("memory_written", crash_after, hard_exit)

    (candidate / "artifacts").mkdir(exist_ok=False)
    _write_bytes(
        candidate / "artifacts" / "probe.xlsx",
        office_payloads["artifacts/probe.xlsx"],
    )
    _boundary("xlsx_written", crash_after, hard_exit)
    _write_bytes(
        candidate / "artifacts" / "probe.pptx",
        office_payloads["artifacts/probe.pptx"],
    )
    _boundary("pptx_written", crash_after, hard_exit)

    files = []
    for name in REQUIRED_FILES:
        path = candidate / name
        files.append(
            {
                "path": name,
                "size": path.stat().st_size,
                "sha256": _sha256(path),
            }
        )
    manifest = {
        "schema_version": PREPARED_SCHEMA,
        "logical_hash": logical_id,
        "physical_id": physical_id,
        "files": files,
    }
    _write_json(candidate / PREPARED, manifest)
    _boundary("prepared_written", crash_after, hard_exit)
    deadline_at = _publication_deadline(deadline_seconds, clock)
    validated = _validate_prepared_until(
        candidate,
        deadline_at=deadline_at,
        clock=clock,
        sleeper=sleeper,
    )
    if validated is None:
        raise StorageIntegrityError(
            "candidate preparation ended as publish_blocked"
        )
    _boundary("prepared_validated", crash_after, hard_exit)

    if commit:
        state = publish_prepared(
            candidate,
            clock=clock,
            sleeper=sleeper,
            _deadline_at=deadline_at,
        )
        if state != "committed":
            raise StorageIntegrityError(
                f"candidate publication ended as {state}"
            )
        _boundary("committed_created", crash_after, hard_exit)
    return candidate


def _safe_manifest_path(candidate, relative):
    relative_path = Path(relative) if isinstance(relative, str) else Path()
    if (
        not isinstance(relative, str)
        or not relative
        or relative_path.is_absolute()
        or ".." in relative_path.parts
        or relative_path.as_posix() != relative
        or relative in {PREPARED, COMMITTED}
    ):
        raise StorageIntegrityError(
            f"unsafe prepared path {relative!r}"
        )
    path = candidate.joinpath(*relative_path.parts)
    if _is_reparse(path) or not path.is_file():
        raise StorageIntegrityError(
            f"prepared member {relative!r} is not a regular file"
        )
    return path


def _is_reparse(path):
    path = Path(path)
    try:
        stat = path.lstat()
    except OSError:
        return False
    reparse_flag = getattr(stat, "st_file_attributes", 0) & 0x400
    return path.is_symlink() or bool(reparse_flag)


def validate_prepared(candidate):
    """Return the manifest only when the candidate is hash-complete."""
    candidate = Path(candidate)
    if _is_reparse(candidate) or not candidate.is_dir():
        raise StorageIntegrityError("candidate is not a regular directory")
    prepared = candidate / PREPARED
    if _is_reparse(prepared) or not prepared.is_file():
        raise StorageIntegrityError("PREPARED.json is absent or irregular")
    try:
        manifest = json.loads(prepared.read_text(encoding="utf-8"))
    except OSError:
        # Callers may retry sharing/access violations without changing state.
        raise
    except (UnicodeError, ValueError) as exc:
        raise StorageIntegrityError(
            f"PREPARED.json cannot be decoded: {exc}"
        ) from exc
    if manifest.get("schema_version") != PREPARED_SCHEMA:
        raise StorageIntegrityError("unsupported prepared format version")
    if manifest.get("logical_hash") != candidate.parent.name:
        raise StorageIntegrityError("logical id does not match its directory")
    if manifest.get("physical_id") != candidate.name:
        raise StorageIntegrityError("physical id does not match its directory")
    entries = manifest.get("files")
    if not isinstance(entries, list):
        raise StorageIntegrityError("prepared files must be a list")
    names = [entry.get("path") for entry in entries if isinstance(entry, dict)]
    if sorted(names) != sorted(REQUIRED_FILES) or len(names) != len(entries):
        raise StorageIntegrityError("prepared member set is incomplete")
    expected_files = set(REQUIRED_FILES) | {PREPARED}
    if (candidate / COMMITTED).exists():
        expected_files.add(COMMITTED)
    actual_files = set()
    actual_directories = set()
    for path in candidate.rglob("*"):
        relative = path.relative_to(candidate).as_posix()
        if _is_reparse(path):
            raise StorageIntegrityError(
                f"candidate member {relative!r} is a reparse point"
            )
        if path.is_dir():
            actual_directories.add(relative)
        elif path.is_file():
            actual_files.add(relative)
        else:
            raise StorageIntegrityError(
                f"candidate member {relative!r} is irregular"
            )
    if actual_directories != {"artifacts"} or actual_files != expected_files:
        raise StorageIntegrityError(
            "candidate contains missing or unexpected members"
        )
    for entry in entries:
        path = _safe_manifest_path(candidate, entry["path"])
        if type(entry.get("size")) is not int or entry["size"] < 0:
            raise StorageIntegrityError("prepared size is invalid")
        if path.stat().st_size != entry["size"]:
            raise StorageIntegrityError(
                f"size mismatch for {entry['path']}"
            )
        digest = entry.get("sha256")
        if (
            not isinstance(digest, str)
            or not re.fullmatch(r"[0-9a-f]{64}", digest)
            or _sha256(path) != digest
        ):
            raise StorageIntegrityError(
                f"hash mismatch for {entry['path']}"
            )
    return manifest


def validate_committed(candidate):
    manifest = validate_prepared(candidate)
    marker = Path(candidate) / COMMITTED
    if _is_reparse(marker) or not marker.is_file():
        raise StorageIntegrityError("commit marker is absent or irregular")
    if marker.stat().st_size != 0:
        raise StorageIntegrityError("commit marker must be empty")
    return manifest


def create_commit_marker(candidate):
    candidate = Path(candidate)
    marker = candidate / COMMITTED
    descriptor = os.open(
        str(marker),
        os.O_CREAT | os.O_EXCL | os.O_WRONLY,
        0o600,
    )
    try:
        os.fsync(descriptor)
    finally:
        os.close(descriptor)
    _sync_parent(candidate)
    return marker


def classify_candidate(candidate):
    candidate = Path(candidate)
    try:
        if (candidate / COMMITTED).exists():
            validate_committed(candidate)
            return "committed"
        validate_prepared(candidate)
        return "prepared"
    except StorageIntegrityError:
        if (candidate / COMMITTED).exists():
            return "corrupt_committed"
        if (candidate / PREPARED).exists():
            return "corrupt_prepared"
        return "abandoned"


def is_retryable_filesystem_error(exc):
    return (
        getattr(exc, "winerror", None) in _RETRYABLE_WINERRORS
        or getattr(exc, "errno", None) in _RETRYABLE_ERRNOS
    )


def _publication_deadline(deadline_seconds, clock):
    if isinstance(deadline_seconds, bool):
        raise ValueError("deadline_seconds must be a finite nonnegative number")
    try:
        duration = float(deadline_seconds)
    except (TypeError, ValueError) as exc:
        raise ValueError(
            "deadline_seconds must be a finite nonnegative number"
        ) from exc
    if not math.isfinite(duration) or duration < 0:
        raise ValueError(
            "deadline_seconds must be a finite nonnegative number"
        )
    return clock() + duration


def _wait_for_retry(attempt, deadline_at, clock, sleeper):
    remaining = deadline_at - clock()
    if remaining <= 0:
        return False
    delay = _DELAYS[attempt] if attempt < len(_DELAYS) else 2.0
    sleeper(min(delay, remaining))
    return True


def _validate_prepared_until(
    candidate,
    *,
    deadline_at,
    clock,
    sleeper,
):
    """Retry only transient filesystem failures until one shared deadline."""
    attempt = 0
    while True:
        try:
            return validate_prepared(candidate)
        except OSError as exc:
            if not is_retryable_filesystem_error(exc):
                raise
            if not _wait_for_retry(
                attempt,
                deadline_at,
                clock,
                sleeper,
            ):
                return None
            attempt += 1


def publish_prepared(
    candidate,
    deadline_seconds=30.0,
    clock=time.monotonic,
    sleeper=time.sleep,
    _deadline_at=None,
):
    """Validate and publish, retrying only idempotent filesystem work."""
    deadline_at = (
        _publication_deadline(deadline_seconds, clock)
        if _deadline_at is None
        else _deadline_at
    )
    attempt = 0
    while True:
        try:
            validate_prepared(candidate)
            try:
                create_commit_marker(candidate)
            except FileExistsError:
                validate_committed(candidate)
            validate_committed(candidate)
            return "committed"
        except OSError as exc:
            if not is_retryable_filesystem_error(exc):
                raise
            if not _wait_for_retry(
                attempt,
                deadline_at,
                clock,
                sleeper,
            ):
                return "publish_blocked"
            attempt += 1


def recover_candidate(candidate, deadline_seconds=30.0):
    candidate = Path(candidate)
    if (candidate / PREPARED).exists():
        return publish_prepared(
            candidate,
            deadline_seconds=deadline_seconds,
        )
    return "abandoned"


def _worker(args):
    templates = Path(args.templates)
    payloads = {
        "artifacts/probe.xlsx": (templates / "template.xlsx").read_bytes(),
        "artifacts/probe.pptx": (templates / "template.pptx").read_bytes(),
    }
    prepare_candidate(
        args.root,
        args.logical,
        args.physical,
        payloads,
        crash_after=args.crash_after,
        hard_exit=True,
    )
    return 0


def _held_handle_worker(args):
    if os.name != "nt":
        raise RuntimeError("held-handle worker requires Windows")
    import ctypes
    from ctypes import wintypes

    kernel32 = ctypes.WinDLL("kernel32", use_last_error=True)
    create_file = kernel32.CreateFileW
    create_file.argtypes = (
        wintypes.LPCWSTR,
        wintypes.DWORD,
        wintypes.DWORD,
        wintypes.LPVOID,
        wintypes.DWORD,
        wintypes.DWORD,
        wintypes.HANDLE,
    )
    create_file.restype = wintypes.HANDLE
    close_handle = kernel32.CloseHandle
    close_handle.argtypes = (wintypes.HANDLE,)
    close_handle.restype = wintypes.BOOL
    handle = create_file(
        str(Path(args.path)),
        0x80000000,  # GENERIC_READ
        0,  # deny read/write/delete sharing
        None,
        3,  # OPEN_EXISTING
        0x80,  # FILE_ATTRIBUTE_NORMAL
        None,
    )
    invalid = ctypes.c_void_p(-1).value
    if handle == invalid:
        raise OSError(ctypes.get_last_error(), "CreateFileW failed")
    try:
        sys.stdout.write("READY\n")
        sys.stdout.flush()
        time.sleep(args.seconds)
    finally:
        if not close_handle(handle):
            raise ctypes.WinError(ctypes.get_last_error())
    return 0


def _exercise_held_handle(
    root,
    logical,
    physical,
    office_payloads,
    artifact_name,
    seconds=0.5,
):
    candidate = prepare_candidate(
        root,
        logical,
        physical,
        office_payloads,
        commit=False,
    )
    command = [
        sys.executable,
        "-m",
        "bench.f0_storage",
        "hold-handle",
        "--path",
        str(candidate / "artifacts" / artifact_name),
        "--seconds",
        str(seconds),
    ]
    process = subprocess.Popen(
        command,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
    )
    ready_holder = []

    def read_ready():
        ready_holder.append(process.stdout.readline().strip())

    reader = threading.Thread(target=read_ready, daemon=True)
    reader.start()
    reader.join(timeout=10)
    if reader.is_alive():
        process.terminate()
        try:
            process.communicate(timeout=5)
        except subprocess.TimeoutExpired:
            process.kill()
            process.communicate(timeout=5)
        reader.join(timeout=1)
        raise StorageIntegrityError(
            "held-handle worker did not become ready within ten seconds"
        )
    ready = ready_holder[0] if ready_holder else ""
    if ready != "READY":
        _, stderr = process.communicate(timeout=10)
        raise StorageIntegrityError(
            f"held-handle worker failed to start: {stderr[-500:]}"
        )
    started = time.monotonic()
    state = publish_prepared(candidate, deadline_seconds=30.0)
    elapsed = time.monotonic() - started
    try:
        _, stderr = process.communicate(timeout=10)
    except subprocess.TimeoutExpired as exc:
        process.terminate()
        try:
            process.communicate(timeout=5)
        except subprocess.TimeoutExpired:
            process.kill()
            process.communicate(timeout=5)
        raise StorageIntegrityError(
            "held-handle worker did not exit after publication"
        ) from exc
    if process.returncode != 0:
        raise StorageIntegrityError(
            f"held-handle worker failed: {stderr[-500:]}"
        )
    if state != "committed":
        raise StorageIntegrityError(
            f"held-handle publication ended as {state}"
        )
    if elapsed < max(0.1, seconds * 0.5):
        raise StorageIntegrityError(
            "held Office handle did not exercise the retry path"
        )
    return candidate, elapsed


def run_spike(
    root,
    cycles=200,
    crash_cycles=50,
    held_handle_cycles=0,
):
    """Exercise normal publication and hard process exits."""
    if type(cycles) is not int or cycles < 1:
        raise ValueError("cycles must be a positive integer")
    if (
        type(crash_cycles) is not int
        or crash_cycles < 0
        or crash_cycles > cycles
    ):
        raise ValueError("crash_cycles must be between zero and cycles")
    if (
        type(held_handle_cycles) is not int
        or held_handle_cycles < 0
        or crash_cycles + held_handle_cycles > cycles
    ):
        raise ValueError(
            "held_handle_cycles must fit within the total cycles"
        )
    if held_handle_cycles and os.name != "nt":
        raise RuntimeError("held-handle cycles require Windows")
    root = Path(root)
    root.mkdir(parents=True, exist_ok=False)
    templates = root / "templates"
    payloads = create_office_templates(templates)
    records = []
    normal_cycles = cycles - crash_cycles - held_handle_cycles

    for index in range(normal_cycles):
        logical = hashlib.sha256(
            f"normal:{index}".encode("ascii")
        ).hexdigest()
        physical = str(uuid.uuid4())
        candidate = prepare_candidate(
            root, logical, physical, payloads
        )
        records.append(
            {
                "kind": "normal",
                "boundary": None,
                "state": classify_candidate(candidate),
            }
        )

    for index in range(crash_cycles):
        boundary = BOUNDARIES[index % len(BOUNDARIES)]
        logical = hashlib.sha256(
            f"crash:{index}".encode("ascii")
        ).hexdigest()
        physical = str(uuid.uuid4())
        command = [
            sys.executable,
            "-m",
            "bench.f0_storage",
            "worker",
            "--store-root",
            str(root),
            "--templates",
            str(templates),
            "--logical",
            logical,
            "--physical",
            physical,
            "--crash-after",
            boundary,
        ]
        completed = subprocess.run(
            command,
            check=False,
            capture_output=True,
            text=True,
            timeout=60,
        )
        if completed.returncode != 99:
            raise StorageIntegrityError(
                "fault worker did not stop at the requested boundary: "
                f"{completed.returncode}: {completed.stderr[-500:]}"
            )
        candidate = candidate_path(root, logical, physical)
        before = classify_candidate(candidate)
        after = recover_candidate(candidate)
        replacement = None
        if after == "abandoned":
            replacement = str(uuid.uuid4())
            replacement_candidate = prepare_candidate(
                root,
                logical,
                replacement,
                payloads,
            )
            after = classify_candidate(replacement_candidate)
        if after in {"corrupt_prepared", "corrupt_committed"}:
            raise StorageIntegrityError(
                f"crash recovery produced {after}"
            )
        records.append(
            {
                "kind": "forced_exit",
                "boundary": boundary,
                "state_before": before,
                "state": after,
                "replacement_physical_id": replacement,
            }
        )

    for index in range(held_handle_cycles):
        logical = hashlib.sha256(
            f"held:{index}".encode("ascii")
        ).hexdigest()
        physical = str(uuid.uuid4())
        candidate, elapsed = _exercise_held_handle(
            root,
            logical,
            physical,
            payloads,
            "probe.xlsx" if index % 2 == 0 else "probe.pptx",
        )
        records.append(
            {
                "kind": "held_handle",
                "boundary": "prepared_adoption",
                "state": classify_candidate(candidate),
                "publish_elapsed_seconds": elapsed,
            }
        )

    committed = 0
    abandoned = 0
    invalid_committed = 0
    logical_commits = {}
    for candidate in (root / "attempts").glob("*/*"):
        state = classify_candidate(candidate)
        if state == "committed":
            committed += 1
            logical_commits.setdefault(candidate.parent.name, []).append(
                candidate.name
            )
            try:
                validate_committed(candidate)
            except StorageIntegrityError:
                invalid_committed += 1
        elif state == "abandoned":
            abandoned += 1
        else:
            raise StorageIntegrityError(
                f"unresolved candidate state {state}"
            )
    duplicates = {
        logical: physical
        for logical, physical in logical_commits.items()
        if len(physical) != 1
    }
    return {
        "schema_version": "brick.f0.storage-summary/1",
        "cycles": cycles,
        "forced_exits": crash_cycles,
        "held_handle_cycles": held_handle_cycles,
        "committed": committed,
        "logical_commits": len(logical_commits),
        "abandoned": abandoned,
        "physical_candidates": committed + abandoned,
        "duplicate_valid_candidates": duplicates,
        "invalid_committed": invalid_committed,
        "directory_renames": 0,
        "records": records,
        "passed": (
            len(records) == cycles
            and invalid_committed == 0
            and committed == cycles
            and len(logical_commits) == cycles
            and not duplicates
        ),
    }


def _parse_args(argv=None):
    parser = argparse.ArgumentParser()
    subparsers = parser.add_subparsers(dest="command", required=True)
    worker = subparsers.add_parser("worker")
    worker.add_argument("--store-root", dest="root", required=True)
    worker.add_argument("--templates", required=True)
    worker.add_argument("--logical", required=True)
    worker.add_argument("--physical", required=True)
    worker.add_argument("--crash-after", choices=BOUNDARIES, required=True)

    holder = subparsers.add_parser("hold-handle")
    holder.add_argument("--path", required=True)
    holder.add_argument("--seconds", type=float, default=0.5)

    spike = subparsers.add_parser("spike")
    spike.add_argument("--store-root", dest="root", required=True)
    spike.add_argument("--cycles", type=int, default=200)
    spike.add_argument("--crash-cycles", type=int, default=50)
    spike.add_argument("--held-handle-cycles", type=int, default=10)
    return parser.parse_args(argv)


def main(argv=None):
    args = _parse_args(argv)
    if args.command == "worker":
        return _worker(args)
    if args.command == "hold-handle":
        return _held_handle_worker(args)
    summary = run_spike(
        args.root,
        cycles=args.cycles,
        crash_cycles=args.crash_cycles,
        held_handle_cycles=args.held_handle_cycles,
    )
    sys.stdout.write(
        json.dumps(summary, indent=2, sort_keys=True) + "\n"
    )
    return 0 if summary["passed"] else 1


if __name__ == "__main__":
    raise SystemExit(main())
